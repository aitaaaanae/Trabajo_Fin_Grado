import time

start_time = time.time()

import procesamiento
from procesamiento import obtain_cols
from procesamiento import to_date
from procesamiento import extract_intervals_3000
from procesamiento import any_treatments
from procesamiento import create_treatment_columns
import datetime as dt
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics
from pptx import Presentation
from sklearn.compose import make_column_selector as selector
from pptx.util import Inches, Pt
from pptx.util import Cm
import io
from PIL import Image
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE 
from scipy.stats import skew
from scipy.stats import kurtosis
import squarify
import matplotlib.patches as mpatches
import matplotlib.colors as mcolors
import inflect
import os
import sys
import researchpy
from scipy.stats import shapiro
import re
import researchpy
import Univ_tools
from Univ_tools import univariate_numerical_slide
from Univ_tools import univariate_piechart_slide
from Univ_tools import univariate_barchart_slide
from Univ_tools import plot_tumors
from Univ_tools import sample_string_function
from Univ_tools import comorbidity_univ_plot
from Univ_tools import divide_patients_per_n_tumors
from Univ_tools import plot_univ_drugs
from Univ_tools import plot_univ_surgeries
from Univ_tools import plot_numerical_tumor
import argparse

      
parser = argparse.ArgumentParser(description='Este script creará un análisis univariante de sus datos en formato .csv que se salvará en un PowerPoint.')
parser.add_argument('-a', '--argument1', type=str, help='Ruta o nombre del archivo .csv para su análisis univariante.')
parser.add_argument('-b', '--argument2', type=str, help='Ruta o nombre del archivo donde se desee guardar la presentación .pptx.')
args = parser.parse_args()

# Access the argument values
argument1_value = args.argument1
argument2_value = args.argument2


prs=Presentation()
df = pd.read_csv(argument1_value, sep=';')

df = obtain_cols(df)
df = to_date(df)
df=extract_intervals_3000(df)
df=any_treatments(df)
df=create_treatment_columns(df, r'surgery_[0-9]+')
df=create_treatment_columns(df, r'drug_[0-9]+')

#Start
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Univariate descriptive analysis of breast cancer patients data'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.size = Pt(30)
tf = body_shape.text_frame
tf.paragraphs[0].font.size = Pt(25)
tf.text = 'Analysis performed for: ' + str(df.shape[0]) 

p = tf.add_paragraph()
p.text = 'Comorbidity analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Demography analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Ginecological analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Tumor type analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Stage analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Time intervals analysis'
p.level = 1

comorbidity_univ_plot(df=df, prs=prs,
                 comorbidities=['autoimmune disease',
 'cardiac insufficiency',
 'diabetes',
 'dislipemia',
 'ex-smoker',
 'gastrointestinal disease',
 'hta',
 'insomnia',
 'ischemic cardiopathology',
 'liver disease',
 'lung disease',
 'musculoskeletal disease',
 'other cardiopathology',
 'psychiatric disorder',
 'renal disease',
 'smoker',
 'thyroid disease',
 'transplant'])  

univariate_numerical_slide(df, prs=prs, sample_string=sample_string_function(df, 
                                                                           col="age_at_diagnosis", mode="normal"), 
                                                                           title_slide="Demography analysis: age at diagnosis", 
                                                                           col="age_at_diagnosis", 
                                                                           xlim=100)

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Gynecological analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Menarche age'
p.level = 1
p = tf.add_paragraph()
p.text = 'Menopause rate'
p.level = 1
p = tf.add_paragraph()
p.text = 'Number of births'
p.level = 1
p = tf.add_paragraph()

p.text = 'Number of abortions'
p.level = 1

p = tf.add_paragraph()

p.text = 'Number of pregnancies'
p.level = 1

univariate_numerical_slide(df, prs=prs, title_slide="Ginecological analysis: menarche age", 
                          sample_string=sample_string_function(df, col="menarche_age", mode="normal"), 
                          col="menarche_age", xlim=20, ylim=700)


univariate_piechart_slide(df, prs=prs, title_slide="Ginecological analysis: menopause at diagnosis", 
                          sample_string=sample_string_function(df, col="menopause", mode="normal"), 
                          col="menopause", colors={"yes":"lightpink", "no":"lightblue"})


univariate_barchart_slide(df, prs=prs,  title_slide="Ginecological analysis: natural births",
                 sample_string=sample_string_function(df, mode="normal", col="natural_birth"),
                     col="natural_birth", regex=None)



univariate_barchart_slide(df, prs=prs, title_slide="Ginecological analysis: pregnancies", 
                          col="pregnancy", sample_string=sample_string_function(df, mode="normal", col="pregnancy"), 
                          regex=None)

univariate_barchart_slide(df, prs=prs, 
                          title_slide="Ginecological analysis: abortions",
                            sample_string=sample_string_function(df,mode="normal", col= "abort"),
col="abort", regex=None)


univariate_barchart_slide(df, prs=prs, title_slide="Ginecological analysis: caesareans", 
                sample_string=sample_string_function(df, mode="normal", col="caesarean"), 
                 col="caesarean", regex=None)



bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Nº of tumors'
p.level = 1

p = tf.add_paragraph()

p.text = 'Histological type'
p.level = 1

p = tf.add_paragraph()
p.text = 'Behavior (in situ vs. invasive, associated in situ)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Staging (at diagnosis and after neoadjuvant therapy) & Neoadjuvance'
p.level = 1


p = tf.add_paragraph()
p.text = 'Celular markers'
p.level = 1

p = tf.add_paragraph()
p.text = 'Clinical subtype'
p.level = 1

df_filtered_1 = divide_patients_per_n_tumors(df)[0]
df_filtered_2 = divide_patients_per_n_tumors(df)[1]


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed (except num. of tumors) are divided according to 3 groups of patients:'

p = tf.add_paragraph()
p.text = 'Group 1: Patients with a single tumor + patients with relapse'+" ("+str(len(df_filtered_1))+ " patients)"
p.level = 1

p = tf.add_paragraph()
p.text = 'Group 2: Patients with more than one tumor in the same date'+" ("+str(len(df_filtered_2))+ " patients)"
p.level = 1



my_palette = sns.color_palette('Set2', 5)
colors = {1: my_palette[0], 2: my_palette[1], 3: my_palette[2], 4:my_palette[3]}
univariate_piechart_slide(df, prs=prs, title_slide="Tumor analysis: number of tumors", 
                          sample_string=sample_string_function(df, mode="normal", col="num_tumors"), 
                          col="num_tumors", regex=None, colors=colors)


univariate_barchart_slide( df, prs=prs, title_slide="Tumor analyisis: number of tumors", 
                          sample_string=sample_string_function(df, mode="normal", col="num_tumors"), 
                          col="num_tumors", regex=None)

plot_tumors(df, title_slide=" histological type ", prs=prs, 
            regex="histological_type_", color_dict={"Ductal": my_palette[0], "Lobular": my_palette[1], "No specific type": my_palette[2]}, 
            exclude="Exclude", explain=" due to ambiguous values of columns Ductal, Lobular, and No specific (treated as unknown)")

plot_tumors(df, title_slide=" behavior ",  prs=prs, 
            regex="behavior_", color_dict={"Invasive": my_palette[0], "In situ": my_palette[1]},
             exclude="Exclude", explain=" due to ambiguous values of columns Invasive and In Situ columns (treated as unknown)")


plot_tumors(df, title_slide=" associated in situ ", prs=prs, 
            regex="associated_col_", color_dict={"Associated in situ": my_palette[0], "No associated in situ": my_palette[1]}, 
             exclude="Exclude", explain=" corresponding to in situ tumors")


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Staging: based on parameters T (tumor size), N (nodal status) and M (metastasis) '
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.size = Pt(30)
tf = body_shape.text_frame
tf.left = Cm(1.27)
tf.top= Cm(5.91)
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Stage at diagnosis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Neoadjuvance therapy in tumors'
p.level = 1
p = tf.add_paragraph()

p = tf.add_paragraph()
p.text = 'Stage after neoadjuvant (only evaluated for patients with neoadjuvant tumors)'
p.level = 1
p = tf.add_paragraph()

plot_tumors(df, title_slide=" stage at diagnosis ", prs=prs, 
             regex="stage_diagnosis_summ_", color_dict = {"0": my_palette[0], "I": my_palette[1], "II": my_palette[2], "III": my_palette[3], "IV": my_palette[4]}, 
             pie="yes", exclude="Exclude", explain=" due to T, N or M parameters not known (treated as unknown)")

plot_tumors(df, title_slide=" neoadjuvance ", prs=prs,  regex="neoadjuvant_", color_dict={"yes":"lightpink", "no":"lightblue"}, 
             pie="yes")

plot_tumors(df, title_slide=" stage after neoadjuvant therapy ", regex="stage_after_neo_summ_", prs=prs, 
            color_dict={"0": my_palette[0], "I": my_palette[1], "II": my_palette[2], "III": my_palette[3], "IV": my_palette[4]}, 
            pie="yes", exclude="x", explain=" due to unknown values of T, N or M parameters", neo="yes")



bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor type analysis: celular markers'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Ki67 percentage per tumor and total'
p.level = 1
p = tf.add_paragraph()
p.text = 'PR per tumor and total'
p.level = 1
p = tf.add_paragraph()
p.text = 'ER per tumor and total'
p.level = 1
p = tf.add_paragraph()

p.text = 'HER2 per tumor and total'
p.level = 1

plot_numerical_tumor(df, "Tumor analysis: ki67 marker ",  prs, xlim=100, 
    regex="ki67_", explain=None, neo=None)



colors_dict = {"Yes": my_palette[0], "No": my_palette[1]}
plot_tumors(df, title_slide="(markers): ER marker", prs=prs, regex="er_positive_cat_",  color_dict=colors_dict)
plot_tumors(df, title_slide="(markers): PR marker", prs=prs,   regex="pr_positive_cat_",  color_dict=colors_dict)
plot_tumors(df, title_slide="(markers): HER-2 marker",  prs=prs, regex="her2_positive_cat_",  color_dict=colors_dict)

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Clinical subtype of tumor: based on variables ER (estrogen receptor), PR (progesterone receptor) and HER2 oncogen'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.size= Pt(30)

tf = body_shape.text_frame
tf.text = 'Types are:'
p = tf.add_paragraph()
p.text = "ER or PR + and HER2 + --> PP"
p.level = 1
p = tf.add_paragraph()
p.text = 'ER or PR + and HER2 - --> PN'
p.level = 1
p = tf.add_paragraph()
p.text = 'ER and PR - and HER2 + --> NP'
p.level = 1
p = tf.add_paragraph()
p.text = 'ER and PR - and HER2 - --> NN'
p.level = 1


colors_dict = {"PN": my_palette[0], "NP": my_palette[1], "NN": my_palette[2], "PP": my_palette[3]}
plot_tumors(df, prs=prs, title_slide=" clinical subtype", regex="type_", color_dict=colors_dict, 
            pie="yes")


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Treatment analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Patients treated by each treatment (hormonetherapy, surgery, radiotherapy)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Types of surgeries'
p.level = 1

p = tf.add_paragraph()
p.text = 'Types of drugs used in hormonetherapy'
p.level = 1

p = tf.add_paragraph()
p.text = 'Number of each type of treatment (surgery, radiotherapy, hormonetherapy)'
p.level = 1
p = tf.add_paragraph()

univariate_piechart_slide(df, prs=prs, col="Any_surgery",
                           sample_string="- Sample size: "+str(len(df))+"\n"
                           "Missing perc of data: 0.0%", 
                           title_slide="Treatment analysis: surgery rate in patients", colors={"yes":"lightpink", "no":"lightblue"})
univariate_piechart_slide(df, prs=prs, col="Any_radiotherapy",
                           sample_string="- Sample size: "+str(len(df))+"\n"
                           "Missing perc of data: 0.0%", 
                           title_slide="Treatment analysis: Rx rate in patients", colors={"yes":"lightpink", "no":"lightblue"})
univariate_piechart_slide(df, prs=prs, col="Any_hormonetherapy",
                           sample_string="- Sample size: "+str(len(df))+"\n"
                           "Missing perc of data: 0.0%", 
                           title_slide="Treatment analysis: Ht rate in patients", colors={"yes":"lightpink", "no":"lightblue"})

plot_univ_surgeries(df, prs)
plot_univ_drugs(df, prs)

univariate_barchart_slide(df, prs=prs, col="N_surgeries",
                           sample_string=sample_string_function(df, mode="normal", col="N_surgeries"), 
                           title_slide="Treatment analysis: nº of surgeries in patients")

univariate_barchart_slide(df, prs=prs, col="n_radio",
                           sample_string=sample_string_function(df, mode="normal", col="n_radio"), 
                           title_slide="Treatment analysis: nº of Rx in patients")

univariate_barchart_slide(df, prs=prs, col="N_hormonetherapies",
                           sample_string=sample_string_function(df, mode="normal", col="N_hormonetherapies"), 
                           title_slide="Treatment analysis: nº of Ht in patients")


bullet_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title

title_shape.top=Cm(1.37)
title_shape.left=Cm(0.45)
title_shape.height=Cm(4.53)
title_shape.width=Cm(24.79)
title_shape.text = 'Time interval analysis'
title_shape.text_frame.paragraphs[0]
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].fontsize = Pt(48)

shape = slide.shapes.add_textbox(top=Cm(6.73), left=Cm(1.47), height=Cm(4.1), width=Cm(22.15))
tf = shape.text_frame		
tf.text = "Variables analyzed:" 

p = tf.add_paragraph()
p.text = '- Time from dx to first treatment (for all patients)'
p.level = 1


p = tf.add_paragraph()
p.text = '- Time from dx first surgery (only patients with no neoadjuvant tumors)'
p.level = 1

p = tf.add_paragraph()
p.text = '- Months between diagnosis of first tumor and recurrence (only patients with recurrence)'
p.level = 1


plot_numerical_tumor(df, title_slide="Time interval analysis: time from dx to treatment (all patients)", 
                     prs=prs, xlim=200, 
                     regex="Time_dx_first_treatment_")

sample_string=sample_string_function(df, mode="tumor_integrated_no_neo", col=None, regex="Time_dx_surgery_no_neo_days_", exclude=200)
univariate_numerical_slide(df, prs=prs, title_slide="Time interval analysis: time from dx to surgery (days) for non neoadjuvant tumors", 
                           sample_string=sample_string,
                  col=None, regex="Time_dx_surgery_no_neo_days_", xlim=200)



univariate_piechart_slide(df, prs=prs,  title_slide="Time interval analysis: recurrence rate",
                 sample_string=sample_string_function(df, mode="normal", col="recurrence"), 
                 col="recurrence", colors={"yes": my_palette[0], "no": my_palette[1]})


sample_string="Sample size: "+str(len(df[df["recurrence_year"].notnull()]))+",\n Missing values: 0.00%"
univariate_numerical_slide(df[df["recurrence"] == "yes"], prs=prs, 
                            title_slide="Time interval analysis: time to recurrence in patients", 
                            sample_string=sample_string,
                  col="Time_to_recurrence_months", regex=None, xlim=200)


prs.save(argument2_value+".pptx")


end_time = time.time()
execution_time = end_time - start_time

print("Script executed in {:.2f} seconds.".format(execution_time))

