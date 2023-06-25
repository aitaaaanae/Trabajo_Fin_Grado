import datetime as dt
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics
from pptx import Presentation
from sklearn.compose import make_column_selector as selector
from pptx.util import Inches, Pt
from pptx.util import Cm
import io
from PIL import Image
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE 
from scipy.stats import skew
from scipy.stats import kurtosis
import squarify
import matplotlib.patches as mpatches
import matplotlib.colors as mcolors
import inflect
import os
import sys
import researchpy
from scipy.stats import shapiro
import re
import researchpy
import matplotlib.ticker as mtick



def sample_string_function(df, mode,  col=None, regex=None, exclude=None, explain=None, normal=None, tumor_filt_num=None):
    df_copy = df.copy()
    if mode=="normal": 
        sample_size = len(df_copy[col])
        missing_perc = ((df_copy[col].isnull().sum())/sample_size)*100
        sample_string= ("- Sample size: {:.2f},\n"
                   "- Missing perc of data: {:.2f}%.")
        sample_string = sample_string.format(sample_size, missing_perc, exclude=exclude, explain=explain)

        if exclude:

            excluded_ehr_col=None
            inconsistencies = 0
            if isinstance(exclude, str):
                excluded = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True).sum()
                mask = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True)
                df_copy.loc[mask, col] = np.nan
            else:
                col_data = df_copy[col]
                excluded_rows = col_data > exclude
                excluded_ehr_col = df_copy.loc[excluded_rows, 'ehr'].tolist()
                inconsistencies = col_data[col_data < 0].count()
                excluded =col_data[col_data > exclude].count()
                                
                col_data[col_data < 0] = np.nan
                col_data[col_data > exclude] = np.nan

                                
            sample_size = len(df_copy[col])
            missing_perc = ((df_copy[col].isnull().sum())/sample_size)*100
            if inconsistencies == 0:  
                sample_string= ("-Sample size: {:.2f},\n"
                   "- Missing perc of data: {:.2f}%,\n"
                   "- Data excluded in the analysis {excluded} (treated as unknown) due to {explain}, \n")
                  # "- Inconsistencies (data < 0) {inconsistencies}" )
                sample_string = sample_string.format(sample_size, missing_perc, excluded=excluded, explain=explain)
            else: 
                sample_string= ("- Sample size: {:.2f},\n"
                   "- Missing perc of data: {:.2f}%,\n"
                   "- Data excluded in the analysis {excluded} (treated as unknown) due to {explain},"
                   "- Inconsistencies (data < 0): {inconsistencies}")
                sample_string = sample_string.format(sample_size, missing_perc, excluded=excluded,
                                                      explain=explain, inconsistencies = inconsistencies)

    if mode == "tumor_filt_num":
        if not exclude: 
            patient_size = len(df_copy[col])
            tumor_sample = df_copy[df_copy["num_tumors"] >= tumor_filt_num]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
            ordinal = get_ordinal(tumor_filt_num)
            sample_string = "- Patients with {ordinal} tumor: {:.2f} ,\n- Of which data for col  {col} is missing for {:.2f}% of tumors."
            sample_string = sample_string.format(tumor_sample_size_patients, missing_perc_tumor, col=col, ordinal=ordinal)
                
        else: 
            excluded_ehr_col=None
            inconsistencies = 0

            if isinstance(exclude, str):
                excluded = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True).sum()
                mask = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True)
                df_copy.loc[mask, col] = np.nan
            else:
                excluded = 0
                excluded_ehr = []
                col_data = df_copy[col]
                excluded_rows = col_data > exclude
                excluded_ehr_col = df_copy.loc[excluded_rows, 'ehr'].tolist()
                inconsistencies = col_data[col_data < 0].count()
                excluded =  col_data[col_data > exclude].count()
                                
                col_data[col_data < 0] = np.nan
                col_data[col_data > exclude] = np.nan

            patient_size = len(df_copy[col])
            tumor_sample = df_copy[df_copy["num_tumors"] >= tumor_filt_num]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
            ordinal = get_ordinal(tumor_filt_num)

            if inconsistencies == 0:  

                sample_string = "- Patients with {ordinal} tumor: {:.2f},\n- Of which data for col  {col} is missing for {:.2f}% of tumors,\n- Data excluded for {excluded} (treated as unknown) due to {explain}. "
                sample_string = sample_string.format(tumor_sample_size_patients,
                                                        missing_perc_tumor, col=col, excluded=excluded, explain=explain, ordinal=ordinal)
            else: 
            
                sample_string = "- Patients with {ordinal} tumor: {:.2f},\n- Of which data for col  {col} is missing for {:.2f}% of tumors,\n- Data excluded for {excluded} (treated as unknown) due to {explain},\n- EHR for excluded patients: {excluded_ehr_col}, \n Inconsistencies (data < 0 ): {inconsistencies}" 
                sample_string = sample_string.format(tumor_sample_size_patients, 
                                                    missing_perc_tumor, col=col, 
                                                     excluded=excluded, explain=explain, excluded_ehr_col=excluded_ehr_col, ordinal=ordinal, inconsistencies=inconsistencies)
    

    if mode == "tumor_filt_num_neoadjuvant":
        if not exclude: 
            patient_size = len(df_copy[col])
            #missing_perc = ((df_copy[col].isnull().sum())/patient_size)*100
            tumor_sample = df_copy[(df_copy["num_tumors"] >= tumor_filt_num) & (df_copy["neoadjuvant_"+str(tumor_filt_num)] == "yes")]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
            ordinal=get_ordinal(tumor_filt_num)
            patient_size = len(df_copy[col])
    #missing_perc = ((df_copy[col].isnull().sum())/patient_size)*100
            tumor_sample = df_copy[(df_copy["num_tumors"] >= tumor_filt_num) & (df_copy["neoadjuvant_"+str(tumor_filt_num)] == "yes")]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
            sample_string= (#"Patients analysed (patient sample size): {:.2f},\n"
                        "- Patients with {ordinal} tumor treated by neoadjuvant: {:.2f},\n"
                        "- Of which data for col  {col} is missing for {:.2f}% of tumors.")
            sample_string = sample_string.format(tumor_sample_size_patients, missing_perc_tumor, col=col, ordinal=ordinal)

        else:
            excluded_ehr_col=None
            inconsistencies = 0
            if isinstance(exclude, str):
                excluded = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True).sum()
                mask = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True)
                df_copy.loc[mask, col] = np.nan
            else:
                excluded = 0
                excluded_ehr = []
                col_data = df_copy[col]
                excluded_rows = col_data > exclude
                excluded_ehr_col = df_copy.loc[excluded_rows, 'ehr'].tolist()
                inconsistencies = col_data[col_data < 0].count()
                excluded =  col_data[col_data > exclude].count()
                                
                col_data[col_data < 0] = np.nan
                col_data[col_data > exclude] = np.nan

                

            patient_size = len(df_copy[col])
            #missing_perc = ((df_copy[col].isnull().sum())/patient_size)*100
            tumor_sample = df_copy[(df_copy["num_tumors"] >= tumor_filt_num) & (df_copy["neoadjuvant_"+str(tumor_filt_num)] == "yes")]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
            ordinal=get_ordinal(tumor_filt_num)

            if  inconsistencies == 0:  
                sample_string= ( "- Patients with {ordinal} tumor treated by neoadjuvant: {:.2f},\n"
                    "- Of which data for col  {col} is missing for {:.2f}% of tumors,\n"
                    "- Data excluded for {excluded} (treated as unknown) due to {explain}.")
                sample_string = sample_string.format(tumor_sample_size_patients, missing_perc_tumor, col=col, excluded=excluded, ordinal=ordinal, explain=explain)
            else: 
                sample_string= ( "- Patients {ordinal} tumor treated by neoadjuvant: {:.2f},\n"
                        "- Of which data for col  {col} is missing for {:.2f}% of tumors,\n"
                        "- Data excluded for {excluded} (treated as unknown) due to {explain},\n"
                        "- Inconsistencies (data < 0): {inconsistencies}")
                sample_string = sample_string.format( tumor_sample_size_patients,
                                                        missing_perc_tumor, col=col,
                                                         excluded=excluded, explain=explain, excluded_ehr_col=excluded_ehr_col, 
                                                         ordinal=ordinal, inconsistencies=inconsistencies)

    if mode == "tumor_filt_num_no_neoadjuvant":
            if not exclude: 
                patient_size = len(df_copy[col])
                tumor_sample = df_copy[(df_copy["num_tumors"] >= tumor_filt_num) & (df_copy["neoadjuvant_"+str(tumor_filt_num)] == "no")]
                tumor_sample_size_patients = len(tumor_sample)
                tumor_sample_size = tumor_sample["num_tumors"].sum()
                missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
                ordinal=get_ordinal(tumor_filt_num)
                patient_size = len(df_copy[col])
               
                tumor_sample_size_patients = len(tumor_sample)
                tumor_sample_size = tumor_sample["num_tumors"].sum()
                missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
                sample_string= (#"Patients analysed (patient sample size): {:.2f},\n"
                            "- Patients with {ordinal} tumor not treated by neoadjuvant: {:.2f},\n"
                            "- Of which data for col  {col} is missing for {:.2f}% of tumors.")
                sample_string = sample_string.format(tumor_sample_size_patients, missing_perc_tumor, col=col, ordinal=ordinal)

            else:
                excluded_ehr=None
                inconsistencies=0
                if isinstance(exclude, str):
                    excluded = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True).sum()
                    mask = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True)
                    df_copy.loc[mask, col] = np.nan
                else:
                   
                    excluded_ehr = []
                    col_data = df_copy[col]
                    excluded_rows = col_data > exclude
                    excluded_ehr = df_copy.loc[excluded_rows, 'ehr'].tolist()
                    excluded = col_data[col_data > exclude].count()
                    inconsistencies = col_data[col_data <0].count()


                    col_data[col_data > exclude] = np.nan
                    col_data[col_data < 0] = np.nan

                patient_size = len(df_copy[col])
                #missing_perc = ((df_copy[col].isnull().sum())/patient_size)*100
                tumor_sample = df_copy[(df_copy["num_tumors"] >= tumor_filt_num) & (df_copy["neoadjuvant_"+str(tumor_filt_num)] == "no")]
                tumor_sample_size_patients = len(tumor_sample)
                tumor_sample_size = tumor_sample["num_tumors"].sum()
                missing_perc_tumor = (tumor_sample[col].isnull().sum()/tumor_sample_size_patients)*100
                
                ordinal=get_ordinal(tumor_filt_num)

                if inconsistencies == 0: 
                    sample_string= ( "- Patients with {ordinal} tumor not treated by neoadjuvant: {:.2f},\n"
                        "- Of which data for col  {col} is missing for {:.2f}% of tumors,\n"
                        "- Data excluded for {excluded} (treated as unknown) due to {explain}.")
                    sample_string = sample_string.format(tumor_sample_size_patients, missing_perc_tumor, 
                                                         col=col, excluded=excluded, ordinal=ordinal, explain=explain)
                else: 
                    sample_string= ( "- Patients {ordinal} tumor not treated by neoadjuvant: {:.2f},\n"
                            "- Of which data for col  {col} is missing for {:.2f}% of tumors,\n"
                            "- Data excluded for {excluded} (treated as unknown) due to {explain},\n"
                            "- Inconsistencies (data < 0): {inconsistencies}. ")
                           # "- EHR of patients excluded: {excluded_ehr_col}")
                    sample_string = sample_string.format( tumor_sample_size_patients,
                                                            missing_perc_tumor, col=col,
                                                            excluded=excluded, explain=explain,
                                                              excluded_ehr_col=excluded_ehr, ordinal=ordinal, inconsistencies=inconsistencies)

           

    
    
    if mode == "tumor_integrated":
    
        patient_size = len(df_copy)
        tumor_sample = df_copy[df_copy["num_tumors"].notnull()]
        tumor_sample_size_patients = len(tumor_sample)
        tumor_sample_size = tumor_sample["num_tumors"].sum()
        missing_tumor = []
        for col_name in df_copy.columns:
            if re.match(regex, col_name):
        
                num_tumor_val = int(col_name.split('_')[-1])
                neo_col = "neoadjuvant_"+str(num_tumor_val)
                col_missing = df_copy.loc[df_copy["num_tumors"] >= num_tumor_val, col_name].isnull().sum()
                missing_tumor.append(col_missing)

        total_missing =sum(missing_tumor)
        total_missing_perc = (total_missing/tumor_sample_size)*100
        sample_string = ("- Patients analysed (patient sample size): {:.2f},\n"
                    "- With a total of {:.2f} tumors (tumor sample size),\n"
                    "- Of which data for cols  {regex} is missing for {:.2f}% of tumors.")
        sample_string = sample_string.format(patient_size, tumor_sample_size, total_missing_perc, regex=regex )

        if exclude:
            excluded_ehr=None
            inconsistencies=0
            
            if isinstance(exclude, str):
                excluded =[]
                for col_name in df_copy.columns: 
                    if re.match(regex, col_name): 
                        excluded.append(len(df_copy[df_copy[col_name].astype(str).str.contains(str(exclude), na=False, case=False)]))
                
                excluded=sum(excluded)
                cols_to_exclude = df_copy.filter(regex=regex).columns
                df_copy.loc[:, cols_to_exclude] = df_copy.loc[:, cols_to_exclude].apply(lambda x: x.astype(str).str.replace(exclude, lambda match: pd.NA, case=False))

            else:
                excluded_ehr = []
                regex_cols = [col for col in df.columns if re.match(regex, col)]
                mask = df_copy[regex_cols] > exclude
                mask_inconsistencies = df_copy[regex_cols] < 0
                excluded = mask.sum().sum()
                inconsistencies = mask_inconsistencies.sum().sum()

                for col_name in df_copy.columns:
                    if re.match(regex, col_name): 
                        col_data = df_copy[col_name]
                        excluded_rows = col_data > exclude
                        excluded_ehr.append(df_copy.loc[excluded_rows, 'ehr'].tolist())
                        col_data[col_data > exclude] = np.nan
                        col_data[col_data < 0] = np.nan
                
            patient_size=len(df_copy)
            tumor_sample = df_copy[df_copy["num_tumors"].notnull()]
            tumor_sample_size_patients = len(tumor_sample)
            tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_tumor = []


            for col_name in df_copy.columns:
                if re.match(regex, col_name):
                    num_tumor_val = int(col_name.split('_')[-1])
                    col_missing = df_copy.loc[df_copy["num_tumors"] >= num_tumor_val, col_name].isnull().sum()
                    missing_tumor.append(col_missing)

            total_missing =sum(missing_tumor)
            total_missing_perc = (total_missing/tumor_sample_size)*100
            
            if inconsistencies ==0: 
                    
                sample_string = ("- Patients analysed (patient sample size): {:.2f},\n"
                        "- With a total of {:.2f} tumors (tumor sample size),\n"
                        "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                        "- Data excluded for {excluded} due to {explain}.")
                sample_string = sample_string.format(patient_size, tumor_sample_size, total_missing_perc, regex=regex, excluded=excluded, explain=explain)
            else: 
                sample_string = ("- Patients analysed (patient sample size): {:.2f},\n"
                        "- With a total of {:.2f} tumors (tumor sample size),\n"
                        "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                        "- Data excluded for {excluded} due to {explain},\n"
                        "- Inconsistencies (data < 0) {inconsistencies} .")
                sample_string = sample_string.format(patient_size, tumor_sample_size, 
                                                     total_missing_perc, regex=regex, 
                                                     excluded=excluded, explain=explain, excluded_ehr=excluded_ehr)

        

    if mode == "tumor_integrated_neo":
        if not exclude:  
            neo_pattern = re.compile(r"neoadjuvant_")
            tumor_sample = df_copy[(df_copy["num_tumors"].notnull()) & (df_copy.filter(regex=neo_pattern) == "yes").any(axis=1)]
            tumor_sample_size_patients = len(tumor_sample)
            neo_cols = df_copy.filter(regex=neo_pattern).columns.to_list()
            tumor_sample_size = tumor_sample[neo_cols].eq('yes').sum().sum()
         
            missing_tumor =[]
            neo_missing =[]
            for col_name in df_copy.columns: 
                if re.match(regex, col_name): 
                    num_tumor_val = int(col_name.split('_')[-1])
                    neo_col = "neoadjuvant_"+str(num_tumor_val)
                    col_missing = tumor_sample.loc[(tumor_sample["num_tumors"] >= num_tumor_val) & (df_copy[neo_col] == "yes"), col_name].reset_index(drop=True).isnull().sum()
                    missing_tumor.append(col_missing)

            total_missing =sum(missing_tumor)
            total_missing_perc = (total_missing/tumor_sample_size)*100
            sample_string = (#"Patients analysed (patient sample size): {:.2f},\n"
                        "- Patients with neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                        "- Of which data for cols  {regex} is missing for {:.2f}% of tumors.")
            sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, regex=regex)

        else: 
            excluded_ehr=None
            inconsistencies=0
            if isinstance(exclude, str):
                excluded =[]
                for col_name in df_copy.columns: 
                    if re.match(regex, col_name): 
                        excluded.append(len(df_copy[df_copy[col_name].astype(str).str.contains(str(exclude), na=False, case=False)]))
                
                excluded=sum(excluded)
                cols_to_exclude = df_copy.filter(regex=regex).columns
                df_copy.loc[:, cols_to_exclude] = df_copy.loc[:, cols_to_exclude].apply(lambda x: x.astype(str).str.replace(exclude, lambda match: pd.NA, case=False))

            else:
                
               
                excluded_ehr = []
                regex_cols = [col for col in df.columns if re.match(regex, col)]
                mask = df_copy[regex_cols] > exclude
                mask_inconsistencies = df_copy[regex_cols] < 0

                excluded = mask.sum().sum()
                inconsistencies = mask_inconsistencies.sum().sum()

                for col_name in df_copy.columns:
                    if re.match(regex, col_name): 
                        col_data = df_copy[col_name]
                        excluded_rows = col_data > exclude
                        excluded_rows_inconsistencies = col_data < 0 
                        excluded_ehr.append(df_copy.loc[excluded_rows, 'ehr'].tolist())
                        df_copy.loc[excluded_rows, col_name] = np.nan
                        df_copy.loc[excluded_rows_inconsistencies, col_name] = np.nan 


                excluded_ehr = list(set(tuple(x) for x in excluded_ehr)) 

            neo_pattern = re.compile(r"neoadjuvant_")
            tumor_sample = df_copy[(df_copy["num_tumors"].notnull()) & (df_copy.filter(regex=neo_pattern) == "yes").any(axis=1)]
            tumor_sample_size_patients = len(tumor_sample)
            neo_cols = df_copy.filter(regex=neo_pattern).columns.to_list()
            tumor_sample_size = tumor_sample[neo_cols].eq('yes').sum().sum()
            
            missing_tumor =[]
            neo_missing =[]
            for col_name in df_copy.columns: 
                if re.match(regex, col_name): 
                    num_tumor_val = int(col_name.split('_')[-1])
                    neo_col = "neoadjuvant_"+str(num_tumor_val)
                    col_missing = tumor_sample.loc[(tumor_sample["num_tumors"] >= num_tumor_val) & (df_copy[neo_col] == "yes"), col_name].reset_index(drop=True).isnull().sum()
                    missing_tumor.append(col_missing)

            total_missing =sum(missing_tumor)
            total_missing_perc = (total_missing/tumor_sample_size)*100

            if inconsistencies == 0:
                    
                sample_string = (#"Patients analysed (patient sample size): {:.2f},\n"
                                "- Patients with neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                                "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                                "- Excluded data {excluded} due to {explain}. " )
                sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, excluded=excluded, regex=regex, explain=explain)
            else:
                sample_string = (#"Patients analysed (patient sample size): {:.2f},\n"
                                "- Patients with neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                                "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                                "- Excluded data {excluded} due to {explain}, \n"
                                "- Inconsistencies (data < 0) {inconsistencies}" )
                sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, 
                                                         regex=regex, excluded=excluded, 
                                                        excluded_ehr_col=excluded_ehr, explain=explain, inconsistencies= inconsistencies) 
    if mode == "tumor_integrated_no_neo":
        if not exclude:  
            # create a copy of the original dataframe
            neo_pattern = re.compile(r"neoadjuvant_")
            tumor_sample = df_copy[(df_copy["num_tumors"].notnull()) & (df_copy.filter(regex=neo_pattern) == "no").any(axis=1)]
            tumor_sample_size_patients = len(tumor_sample)
            neo_cols = df_copy.filter(regex=neo_pattern).columns.to_list()
            tumor_sample_size = tumor_sample[neo_cols].eq('no').sum().sum()
            #tumor_sample_size = tumor_sample["num_tumors"].sum()
            missing_tumor =[]
            neo_missing =[]
            for col_name in df_copy.columns: 
                if re.match(regex, col_name): 
                    num_tumor_val = int(col_name.split('_')[-1])
                    neo_col = "neoadjuvant_"+str(num_tumor_val)
                    col_missing = tumor_sample.loc[(tumor_sample["num_tumors"] >= num_tumor_val) & (df_copy[neo_col] == "no"), col_name].reset_index(drop=True).isnull().sum()
                    missing_tumor.append(col_missing)

            total_missing =sum(missing_tumor)
            total_missing_perc = (total_missing/tumor_sample_size)*100
            sample_string = ("- Patients not treated by neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                        "- Of which data for cols  {regex} is missing for {:.2f}% of tumors.")
            sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, regex=regex)

        else: 
            excluded_ehr=None
            inconsistencies = 0
            if isinstance(exclude, str):
                excluded =[]
                for col_name in df_copy.columns: 
                    if re.match(regex, col_name): 
                        excluded.append(len(df_copy[df_copy[col_name].astype(str).str.contains(str(exclude), na=False, case=False)]))
                
                excluded=sum(excluded)
                cols_to_exclude = df_copy.filter(regex=regex).columns
                df_copy.loc[:, cols_to_exclude] = df_copy.loc[:, cols_to_exclude].apply(lambda x: x.astype(str).str.replace(exclude, lambda match: pd.NA, case=False))

            else:
                
               
                excluded_ehr = []
                regex_cols = [col for col in df.columns if re.match(regex, col)]
                mask = df_copy[regex_cols] > exclude
                mask_inconsistencies = df_copy[regex_cols] < 0 

                excluded = mask.sum().sum()

                for col_name in df_copy.columns:
                    if re.match(regex, col_name): 
                        col_data = df_copy[col_name]
                        excluded_rows = col_data > exclude
                        excluded_rows_inconsistencies = col_data < 0
                        excluded_ehr.append(df_copy.loc[excluded_rows, 'ehr'].tolist())
                        df_copy.loc[excluded_rows, col_name] = np.nan 
                        df_copy.loc[excluded_rows_inconsistencies, col_name] = np.nan


                excluded_ehr = list(set(tuple(x) for x in excluded_ehr)) 
            neo_pattern = re.compile(r"neoadjuvant_")
            tumor_sample = df_copy[(df_copy["num_tumors"].notnull()) & (df_copy.filter(regex=neo_pattern) == "no").any(axis=1)]
            tumor_sample_size_patients = len(tumor_sample)
            neo_cols = df_copy.filter(regex=neo_pattern).columns.to_list()
            tumor_sample_size = tumor_sample[neo_cols].eq('no').sum().sum()
            
            missing_tumor =[]
            neo_missing =[]
            for col_name in df_copy.columns: 
                if re.match(regex, col_name): 
                    num_tumor_val = int(col_name.split('_')[-1])
                    neo_col = "neoadjuvant_"+str(num_tumor_val)
                    col_missing = tumor_sample.loc[(tumor_sample["num_tumors"] >= num_tumor_val) & (df_copy[neo_col] == "no"), col_name].reset_index(drop=True).isnull().sum()

                    missing_tumor.append(col_missing)

            total_missing =sum(missing_tumor)
            total_missing_perc = (total_missing/tumor_sample_size)*100

            if inconsistencies == 0:
                    
                sample_string = (#"Patients analysed (patient sample size): {:.2f},\n"
                                "- Patients not treated by neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                                "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                                "- Excluded data {excluded} due to {explain}. " )
                sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, excluded=excluded, regex=regex, explain=explain)
            else:
                sample_string = (#"Patients analysed (patient sample size): {:.2f},\n"
                                "- Patients not treated by neoadjuvant {:.2f}, with a total of {:.2f} tumors (tumor sample size),\n"
                                "- Of which data for cols  {regex} is missing for {:.2f}% of tumors,\n"
                                "- Excluded data {excluded} due to {explain}, \n"
                                "- Inconsistent data (< 0 ): {inconsistencies}" )
                sample_string = sample_string.format(tumor_sample_size_patients, tumor_sample_size, total_missing_perc, 
                                                         regex=regex, excluded=excluded, 
                                                        excluded_ehr_col=excluded_ehr, explain=explain, inconsistencies=inconsistencies) 

    return sample_string
 
def comorbidity_univ_plot(df, prs, comorbidities): 
    excluded_values = (df[comorbidities] == "si,no") | (df[comorbidities] == "no,si")
    excluded_values_count = excluded_values.any(axis=1).sum()


    used_data_comment = "Used data: "
    most_common_comment = "Most common comorbidity: "
    missing_values_comment = "Total missing values for all comorbidities: "
    excluded_values_comment = "Number of excluded rows (yes, no in same cell for any of comorbidities): "

    # Get most common comorbidity
    most_common = df[comorbidities].apply(pd.Series.value_counts).T.sort_values(by=['si'], ascending=False).index[0]
    most_common_comment += f"{most_common}, "

    # Get number of missing values
    missing_values = df[comorbidities].isnull().sum()
    missing_values_comment += f"{missing_values.sum()}, "

    # Get number of excluded values
    excluded_values = (df[comorbidities] == "si,no") | (df[comorbidities] == "no,si")
    excluded_values_count = excluded_values.any(axis=1).sum()
    excluded_values_comment += f"{excluded_values_count},"

    # Get total number of rows in dataset
    sample_size = len(df) - excluded_values_count
    used_data_comment += f"{sample_size}"

    # Join comments into a single string
    comments = f"{used_data_comment}\n{most_common_comment}\n{missing_values_comment}\n{excluded_values_comment}"



    colors = sns.color_palette('Set2', 2)

    # Calculate proportions of si and no for each comorbidity
    si_prop = []
    no_prop = []
    for comorbidity in comorbidities:
        si = df[comorbidity].value_counts().get('si', 0)
        no = df[comorbidity].value_counts().get('no', 0)
        total = len(df[comorbidity].dropna()) - excluded_values[comorbidity].sum()
        si_prop.append(si/total * 100)
        no_prop.append(no/total * 100)

    # Create stacked bar chart
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.bar(comorbidities, si_prop, color=colors[0], label='si')
    ax.bar(comorbidities, no_prop, bottom=si_prop, color=colors[1], label='no')

    # Set y-axis limits to [0, 100]
    ax.set_ylim([0, 100])

    ax.set_title('Comorbidity proportions')
    ax.set_xlabel('Comorbidity')
    ax.set_ylabel('Proportion')
    ax.legend()
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    plt.savefig('comorbidities_stacked_bar_prop.png', bbox_inches='tight')


    img_path = "comorbidities_stacked_bar_prop.png"
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    title = slide.shapes.title
    title.left=Cm(6.95)
    title.top=Cm(1.21)
    title.height=Cm(1.31)
    title.width=Cm(13.69)
    title.text = "Comorbidity analysis"

    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(30)


    left=Cm(2.79)
    top=Cm(2.54)
    pic=slide.shapes.add_picture(img_path, left, top, width=Cm(20.74), height=Cm(12.91))


    left= Cm(0.18)
    top= Cm(15.45)
    width = Cm(24.9)
    height = Cm(3.3)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = comments
    #text 
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(12)	



def univariate_numerical_slide(df, prs, sample_string, title_slide, col=None, regex=None, xlim=None, ylim=None):
    df_copy = df.copy()  # Create a copy of the DataFrame

    if regex:
        cols = [col for col in df_copy.columns if re.match(regex, col)]
        id_cols = "ehr"
        df_copy = pd.melt(df_copy, id_vars=id_cols, value_vars=cols, var_name='variable', value_name='value')
        col = "value"
    
    if xlim: 
        xlim_mask = df_copy[(df_copy[col] < 0) | (df_copy[col] > xlim)]
        df_copy.loc[xlim_mask.index, col] = np.nan


    fig = plt.figure(figsize=(10, 10))
    gs = fig.add_gridspec(2, 2)

    # Add histograms subplot
    ax1 = fig.add_subplot(gs[0, 0])
    sns.histplot(data=df_copy, x=col, color="skyblue", kde=True, alpha=0.5, ax=ax1)
    ax1.set_xlabel(col)
    ax1.set_ylabel("Counts")

    # Perform Shapiro-Wilk test for normality
    if len(df_copy[col].dropna()) > 3: 
        stat, p_value = shapiro(df_copy[col].dropna())
        # Determine whether to display mean and standard deviation or median and IQR
        if p_value >= 0.05:
            measure = "Mean"
            center = statistics.mean(df_copy[col].dropna())
            spread = statistics.stdev(df_copy[col].dropna())
        else:
            measure = "Median"
            center = statistics.median(df_copy[col].dropna())
            spread = np.nanpercentile(df_copy[col].dropna(), 75) - np.nanpercentile(df_copy[col].dropna(), 25)

        # Add measure and center information to histogram
       # ax1.text(0.95, 0.9, f"{measure}: {center:.2f}\nSpread: {spread:.2f}", 
        #        transform=ax1.transAxes, ha="right", va="top")

    # Plot boxplot
    ax2 = fig.add_subplot(gs[0, 1])
    sns.boxplot(data=df_copy, x=col, color="dodgerblue", ax=ax2, orient="h")
    ax2.set_xlabel(col)
    ax2.set_ylabel("")

    # Set darker blue color for density plot
    density_color = "#335B8E"
    sns.histplot(data=df_copy, x=col, kde=True, color=density_color, alpha=0.5, ax=ax1)

    # Styling adjustments
    plt.tight_layout()
    sns.set(style="whitegrid")
    sns.despine()

    # Set x-axis labels based on variable name or regex
    if regex:
        ax2.set_xlabel(regex)
        ax1.set_xlabel(regex)
    else:
        ax2.set_xlabel(col)
        ax1.set_xlabel(col)
    
    if xlim: 
        
        ax1.set_xlim([0,xlim])
        ax2.set_xlim([0, xlim])
    
    if ylim: 
        ax1.set_ylim([0,ylim])

    

    
    if len(df_copy[col].dropna()) > 3: 
        minimum = df_copy[col].dropna().min()
        maximum = df_copy[col].dropna().max()
        range_ = maximum - minimum
        iqr = np.nanpercentile(df_copy[col], 75) - np.nanpercentile(df_copy[col], 25)
        outliers = df_copy[col].dropna()[~df_copy[col].between(
        np.nanpercentile(df_copy[col], 25) - 1.5 * iqr,
        np.nanpercentile(df_copy[col], 75) + 1.5 * iqr)]
        outliers_l = len(outliers)
        # Create text box for statistical summary on histograms
        if p_value >= 0.05:
            ax1.text(0.75, 0.9, f"{measure}: {center:.2f}\nSpread (SD): {spread:.2f}\nMin: {minimum:.2f}\nMax: {maximum:.2f}\nRange: {range_:.2f}",
                    horizontalalignment='right',
                    verticalalignment='top',
                    transform=ax1.transAxes,
                    fontsize=10)

            # Create text box for outlier comment on boxplot
            ax2.text(0.95, 0.9, f"{outliers_l} outliers",
                    horizontalalignment='right',
                    verticalalignment='top',
                    transform=ax2.transAxes,
                    fontsize=10)
        else:
            ax1.text(0.75, 0.9, f"{measure}: {center:.2f}\nSpread (IQR): {spread:.2f}\nMin: {minimum:.2f}\nMax: {maximum:.2f}\nRange: {range_:.2f}",
                    horizontalalignment='right',
                    verticalalignment='top',
                    transform=ax1.transAxes,
                    fontsize=10)

            # Create text box for outlier comment on boxplot
            ax2.text(0.95, 0.9, f"{outliers_l} outliers",
                    horizontalalignment='right',
                    verticalalignment='top',
                    transform=ax2.transAxes,
                    fontsize=10)

    
    # Save the figure
    img_path = f"numerical_variable_univ_{col}.png"
    fig.savefig(img_path, bbox_inches="tight")

    # Create the slide
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    # Set the title
    title = slide.shapes.title
    title.text = title_slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True

    # Add the sample string
    left = Cm(0.75)
    top = Cm(3.32)
    width = Cm(22.29)
    height = Cm(2.9)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.paragraphs[0].text = sample_string
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
    text_frame.word_wrap = True
    text_frame.paragraphs[0].font.size = Pt(17)

    # Add the image
    pic = slide.shapes.add_picture(img_path, left=Cm(2.05), top=Cm(6.76), width=Cm(20.52), height=Cm(10.13))

def divide_patients_per_n_tumors(df):
    df_filtered_1 = df[df["num_tumors"] == 1]
    df_filtered_2 = df[df['num_tumors'] > 1]
    date_columns = df_filtered_2.filter(regex='tumor_date_').columns.tolist()
    df_filtered_2 = df_filtered_2[df_filtered_2[date_columns].nunique(axis=1) == 1]
    df_filtered_3 = df[df['num_tumors'] > 1]
    date_columns = df_filtered_3.filter(regex='tumor_date_').columns.tolist()
    df_filtered_3 = df_filtered_3[df_filtered_3[date_columns].nunique(axis=1) > 1]
        
    df_combined = pd.concat([df_filtered_1, df_filtered_3])
    
    return df_combined, df_filtered_2


def univariate_barchart_slide(df, prs, title_slide, sample_string,  col=None, regex=None, title=None, exclude=None):
    if regex:
        cols = [col for col in df.columns if re.match(regex, col)]
        id_cols = "ehr"
        df = pd.melt(df, id_vars=id_cols, value_vars=cols, var_name='variable', value_name='value')
        col = "value"

    if exclude:
        excluded = len(df[df[col].astype(str).str.contains(str(exclude), na=False, case=False)])
        df = df[~df[col].astype(str).str.contains(str(exclude), na=False, case=False)]

    fig, ax = plt.subplots(figsize=(10, 6))
    df_grouped = df.groupby(col).size()
    df_grouped.plot(kind='bar', rot=0, color="lightskyblue", ax=ax)

    # Set aesthetics
    ax.set_ylabel("")
    ax.set_title(title)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(axis='x', labelrotation=45)

    # Calculate category count and mode
    if df[col].notna().any(): 
        categories = len(df[col].value_counts())
        mode = statistics.mode(df[col].dropna())
    else: 
        mode="No mode for empty data"
        categories="No categories"

    # Add count numbers on top of each bar
    for i, count in enumerate(df_grouped.values):
        ax.text(i, count, str(count), ha='center', va='bottom')

    # Add comment on the graph
    comment = f"Categories: {categories}\nMode: {mode}"
    ax.text(0.95, 0.95, comment, ha='right', va='top', transform=ax.transAxes, fontsize=10)

    fig.savefig("barchart_variable_"+str(col)+".png")
    img_path = "barchart_variable_"+str(col)+".png"
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = title_slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(25)


        
    pic=slide.shapes.add_picture(img_path, left=Cm(3.97), top=Cm(5.97), width=Cm(17.46), height=Cm(11.13))

    #pic=slide.shapes.add_picture("Barchart_treatment.png", left, top, width=Cm(14), height=Cm(11.02))

    left= Cm(1.12)
    top= Cm(3.07)
    width = Cm(22.29)
    height = Cm(2.9)
    gap = Cm(0.5) 
        #shape = slide.shapes.add_textbox(left, top, width, height)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text=sample_string
    font = run.font
    font.size = Pt(14)

def univariate_piechart_slide(df, prs, sample_string, title_slide,  col=None, regex=None, 
                      colors=None, title=None, exclude=None):
    df_copy = df.copy()  # Create a copy of the input DataFrame

    if regex:
        cols = [col for col in df_copy.columns if re.match(regex, col)]
        id_cols = "ehr"

       
        df_copy = pd.melt(df_copy, id_vars=id_cols, value_vars=cols, var_name='variable', value_name='value')
        col = "value"

    if exclude:
        excluded = len(df_copy[df_copy[col].astype(str).str.contains(str(exclude), na=False, case=False)])
        mask = df_copy[col].astype(str).str.contains(exclude, na=False, regex=True)
        df_copy.loc[mask, col] = np.nan

  
    categories = df_copy[col].dropna().unique()

    labels = []
    counts = []

    for category in categories:
        count = len(df_copy[df_copy[col] == category])
        labels.append(category)
        counts.append(count)

    fig, ax = plt.subplots(figsize=(8,6))
    wedges, _, autotexts = ax.pie(
        counts,
        colors=[colors[label] for label in labels],
        autopct=lambda pct: "{:.1f}%\n({:.0f})".format(pct, pct / 100. * sum(counts)),
        textprops={'fontsize': 10, 'fontweight': 'bold'}
    )
    legend = ax.legend(wedges, labels, title=regex if regex else col, loc='center left', bbox_to_anchor=(1, 0.5))
    legend.set_bbox_to_anchor((-0.1, -0.05))

    # Set aesthetics
    ax.axis('equal')
    plt.setp(autotexts, weight='normal')
    plt.title(title, fontweight='bold')

    
    if df_copy[col].notna().any():  
        mode = statistics.mode(df_copy[col].dropna())
       
        categories = len(df_copy[col].value_counts())
    else: 
        mode= "No mode for empty data"
        categories = "No categories"

    # Add comment on the graph
    comment = f"Categories: {categories}\nMode: {mode}"
    ax.text(
        0.95, 0.95, comment,
        ha='right', va='top', transform=ax.transAxes,
        fontsize=12,
        bbox=dict(facecolor='white', edgecolor='gray', boxstyle='round,pad=0.5')
    )

    fig.savefig("piechart_variable_univ_"+str(col)+".png", bbox_inches="tight")
    img_path = "piechart_variable_univ_"+str(col)+".png"

    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = title_slide
    title.left = Cm(1.37)
    title.top = Cm(0.26)
    title.width = Cm(22.66)
    title.height = Cm(0.97)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(23)


    left=Cm(2.4)
    top=Cm(3.82)
    #width=Cm(14.9)
    #height=Cm(11.73)
    pic=slide.shapes.add_picture(img_path, left, top)

    left= Cm(1.37)
    top= Cm(1.57)
    width = Cm(22.29)
    height = Cm(2.9)
    gap = Cm(0.5) 
    #shape = slide.shapes.add_textbox(left, top, width, height)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text=sample_string
    font = run.font
    font.size = Pt(14)
   


def get_ordinal(num):
    """
    Convert number to ordinal form
    """
    p = inflect.engine()
    return p.ordinal(num)

def plot_numerical_tumor(df, title_slide,  prs,  xlim, 
    regex, explain=None, neo=None):

    
    df_1 = divide_patients_per_n_tumors(df)[0]
    df_2 = divide_patients_per_n_tumors(df)[1]
                
    if neo:
        mode1 = "tumor_filt_num_neoadjuvant"
        mode2 = "tumor_filt_num_neoadjuvant"
        mode3 = "tumor_integrated_neo"
        tumor_filt_num_1 = 1
    else:
        mode1 = "normal"
        mode2 = "tumor_filt_num"
        mode3 = "tumor_integrated"
        tumor_filt_num_1 = None 

   #  univariate_numerical_slide(df_filtered_1, title_slide=str(title_slide)+" for group 1 of patients", prs=prs, 
      #             sample_string=sample_string_function(df_filtered_1, mode="normal", col=regex+"1"), 
       #            col=regex+"1", regex=None, xlim=100)

    n_tumors = df.filter(regex="tumor_date_").shape[1]
    for i in range(1, n_tumors+1):
        if regex+str(i) in df_1.columns: 
            if df_1[regex+str(i)].notnull().any:
                sample_string=sample_string_function(df_1, mode="tumor_filt_num" , col=regex+str(i),
                                              regex=None, exclude=None, explain=None, tumor_filt_num=i)
                univariate_numerical_slide(df_1, title_slide=str(title_slide)+" for tumor " +str(i)+"of group 1 of patients", 
                                            col=regex+str(i), prs=prs, 
                         sample_string=sample_string, regex=None, xlim=xlim)
    
    
    sample_string=sample_string_function(df_1, mode="tumor_integrated", col=None, regex=regex)
    univariate_numerical_slide(df_1, 
                            title_slide=str(title_slide)+" for all tumors for patients for group 1 of patients (progression)", prs=prs, 
                            sample_string=sample_string, col=None, regex=regex, xlim=xlim)

         

    sample_string=sample_string_function(df_2, mode="tumor_integrated" , col=None, regex=regex)
    univariate_numerical_slide(df_2, title_slide=str(title_slide)+" for all tumors for patients for group 2 of patients", prs=prs, 
                 sample_string=sample_string, col=None, regex=regex, xlim=xlim)





def plot_tumors(df,  
    title_slide, regex, prs,  
    color_dict, pie=None, bar=None, exclude=None, explain=None, neo=None):

    df_1 = divide_patients_per_n_tumors(df)[0]
    df_2 = divide_patients_per_n_tumors(df)[1]
                
    if neo:
        mode1 = "tumor_filt_num_neoadjuvant"
        mode2 = "tumor_filt_num_neoadjuvant"
        mode3 = "tumor_integrated_neo"
        tumor_filt_num_1 = 1
    else:
        mode1 = "normal"
        mode2 = "tumor_filt_num"
        mode3 = "tumor_integrated"
        tumor_filt_num_1 = None

         

             # Group 4

    n_tumors = df.filter(regex="tumor_date_").shape[1]
    for i in range(1, n_tumors+1):
        if regex+str(i) in df_1.columns:
            if df_1[regex+str(i)].notnull().any():

                univariate_piechart_slide(df_1, prs=prs, 
                                          title_slide="Tumor analysis: "+str(title_slide)+" of tumor "+str(i)+" for patients in group 1",
                            sample_string=sample_string_function(df_1, mode=mode2, col=regex+str(i), tumor_filt_num=i , 
                                                                 exclude=exclude, explain=explain),col=regex+str(i), 
                                colors=color_dict, regex=None,  exclude=exclude)
                
    univariate_piechart_slide(df_1, prs=prs, 
                               title_slide="Tumor analysis: "+str(title_slide)+" for all tumors for patients in group 1",
                                sample_string=sample_string_function(df_1, 
                                                                     col=None, mode=mode3,
                                                                       regex=regex, exclude=exclude, explain=explain), 
                                col=None, regex=regex, exclude=exclude, colors=color_dict)



 
      
            
    univariate_piechart_slide(df_2, prs, title_slide="Tumor analysis: "+str(title_slide)+" for all tumors for patients in group 2",
                                sample_string=sample_string_function(df_2, mode=mode3, col=None, regex=regex, 
                                                                     exclude=exclude, explain=explain), 
                                                                     col=None, regex=regex, exclude=exclude, colors=color_dict)
    


def plot_univ_surgeries(df, prs): 
    
    sample_string = "Sample size (patients treated by surgery): "+str(len(df[df["Any_surgery"] == "Yes"]))+",\n Number of missing values: 0.0%"

    def count_unique_surgeries(row):
        return len(set(row))

    # filter the dataframe to only include columns with regex pattern "surgery_\d+"
    surgery_cols = df.filter(regex="surgery_\d+")

    # get the unique categories in the surgery columns
    surgery_categories = surgery_cols.values.flatten()
    surgery_categories = [x for x in surgery_categories if str(x) != 'nan']
    surgery_categories = list(set(surgery_categories))

    # define a color map for the surgery categories
    cmap = plt.get_cmap("Set3")
    colors = [cmap(i) for i in range(len(surgery_categories))]

    # apply the function to each row of the dataframe, but only for the surgery columns
    counts = surgery_cols.apply(count_unique_surgeries, axis=1).value_counts().sort_index()

    # calculate percentages
    total_patients = len(df)
    percentages = [sum(surgery_cols.apply(lambda x: surgery in x.unique(), axis=1))/total_patients*100 for surgery in surgery_categories]

    # plot a horizontal bar chart
    fig, ax = plt.subplots()
    bars = ax.barh(surgery_categories, percentages, color=colors)
    # set the plot title and axis labels
    plt.title('Percentage of patients treated with each surgery type')
    plt.xlabel('Percentage of patients')
    plt.ylabel('Surgery type')

    # format x-axis ticks as percentages
    ax.xaxis.set_major_formatter(mtick.PercentFormatter())

    # create a list of patches for each category
    patches = [mpatches.Patch(color=colors[i], label=surgery_categories[i]) for i in range(len(surgery_categories))]

    # add the legend to the plot
    plt.legend(handles=patches, title='Surgery type', loc='upper right')

    # remove y-axis tick labels
    plt.yticks([])

    # show the plot
    plt.savefig("surg_perc.png")
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Treatments: patients treated (at least once) by each surgery"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.bold = True

    pic=slide.shapes.add_picture("surg_perc.png", left=Cm(4.33), top=Cm(6.25), width=Cm(15.89), height=Cm(10.13))

        #pic=slide.shapes.add_picture("Barchart_treatment.png", left, top, width=Cm(14), height=Cm(11.02))

    left= Cm(1.12)
    top= Cm(3.07)
    width = Cm(22.29)
    height = Cm(2.9)
    gap = Cm(0.5) 
            #shape = slide.shapes.add_textbox(left, top, width, height)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    text_frame.word_wrap = True # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text=sample_string
    font = run.font
    font.size = Pt(17)
 

def plot_univ_drugs(df, prs): 
    sample_string="Sample size (patients treated by hormonetherapy): "+str(len(df[df["Any_hormonetherapy"] == "Yes"]))+",\n Number of missing values: 0.0%"

    # define a list of columns to search for the string 'acetato de megestrol'
    drug_cols = [col for col in df.columns if col.startswith('drug_')]

    # replace the string 'acetato de megestrol' with 'megestrol acetate' in the selected columns
    df[drug_cols] = df[drug_cols].replace('acetato de megestrol', 'megestrol acetate', regex=True)

    def count_unique_drugs(row):
        return len(set(row))

    # filter the dataframe to only include columns with regex pattern "drug_\d+"
    drug_cols = df.filter(regex="drug_\d+")

    # get the unique categories in the drug columns
    drug_categories = drug_cols.values.flatten()
    drug_categories = [x for x in drug_categories if str(x) != 'nan']
    drug_categories = list(set(drug_categories))

    # apply the function to each row of the dataframe, but only for the drug columns
    counts = drug_cols.apply(count_unique_drugs, axis=1).value_counts().sort_index()

    # calculate the percentage of patients for each drug type
    percentages = pd.DataFrame(columns=drug_categories)
    for drug in drug_categories:
        percentages[drug] = [sum(drug_cols.apply(lambda x: drug in x.unique(), axis=1))/len(drug_cols)*100]
    # plot a horizontal bar chart
    fig, ax = plt.subplots()

    # set up the rainbow color palette
    n = len(drug_categories)
    rainbow = plt.get_cmap('rainbow')
    colors = [rainbow(i/n) for i in range(n)]

    # create the horizontal bar chart with drug names on the y-axis
    bars = ax.barh(drug_categories, percentages.values[0], color=colors)

    # set the plot title and axis labels
    plt.title('Percentage of patients treated with each hormone therapy drug')
    plt.xlabel('Percentage of patients (%)')
    plt.ylabel('Hormone therapy drug')

    # add a legend for the colors
    patches = [mpatches.Patch(color=colors[i], label=drug_categories[i]) for i in range(len(bars))]
    #plt.legend(handles=patches, title='Hormone therapy drug', loc='upper right', ncol=2, fontsize=7)

    new_labels = []
    for label in drug_categories:
        if ' ' in label:
            idx = label.find(' ')
            new_label = label[:idx] + '\n' + label[idx+1:]
            new_labels.append(new_label)
        else:
            new_labels.append(label)
            
    ax.set_yticks(np.arange(len(new_labels)))
    ax.set_yticklabels(new_labels)
    
    plt.yticks(fontsize=8)
    plt.savefig("hormonetherapy_new.png")


    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Treatments: patients treated (at least once) by each drug in hormonetherapy"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(38)
    left=Cm(0.71)
    top=Cm(4.33)

    pic=slide.shapes.add_picture("hormonetherapy_new.png", left=Cm(4.33), top=Cm(6.25), width=Cm(15.89), height=Cm(10.13))

    left= Cm(1.12)
    top= Cm(3.07)
    width = Cm(22.29)
    height = Cm(2.9)
    gap = Cm(0.5) 
            #shape = slide.shapes.add_textbox(left, top, width, height)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.paragraphs[0].font.bold = True
    text_frame.word_wrap = True # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text=sample_string
    font = run.font
    font.size = Pt(17)