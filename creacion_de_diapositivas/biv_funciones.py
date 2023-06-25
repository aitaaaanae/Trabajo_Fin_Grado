import datetime as dt
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics as stats
from pptx import Presentation
from sklearn.compose import make_column_selector as selector
from pptx.util import Inches, Pt
from pptx.util import Cm
import io
import re
from PIL import Image
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE 
import inflect
import os
import sys
import matplotlib.patches as mpatches
from scipy import stats
import matplotlib.pyplot as plt
import matplotlib.pyplot as plt
import numpy as np
import statistics
from scipy.stats import shapiro
import math
from scipy.stats import chi2_contingency, f_oneway, shapiro
from pingouin import kruskal




def get_ordinal(num):
    """
    Convert number to ordinal form
    """
    p = inflect.engine()
    return p.ordinal(num)


def sample_string_function(df, mode, divide_col, col=None, exclude=None, explain=None, tumor_filt_num=None):
    sample_dict = {}
    unique_categories = df[divide_col].dropna().unique()
    filtered_dfs = {}

    for category in unique_categories:
        filtered_dfs[category] = df[df[divide_col] == category].copy()

    if mode == "normal":
        if exclude:
            for category, df_category in filtered_dfs.items():
                if df_category[col].notnull().any():
                    excluded_ehr = None
                    if isinstance(exclude, str):
                        excluded = df_category[col].str.contains(exclude, na=False, regex=True).sum()
                        mask = df_category[col].str.contains(exclude, na=False, regex=True)
                        df_category.loc[mask, col] = np.nan
                    else:
                        col_data = df_category[col].copy()
                        excluded_rows = col_data > exclude
                        excluded_ehr = df_category.loc[excluded_rows, 'ehr'].tolist()
                        excluded = col_data[col_data > exclude].count()
                        df_category.loc[df_category[col] > exclude, col] = np.nan

                    sample_size = len(df_category[col])
                    missing_perc = (df_category[col].isnull().sum() / sample_size) * 100
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc,
                        "excluded": excluded,
                        "excluded_ehr": excluded_ehr
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0,
                        "excluded": 0,
                        "excluded_ehr": 0
                    }

            text_output = ""
            for category, data in sample_dict.items():
                
                text_output += f"{category}: Sample size: {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%. Data excluded: {data['excluded']} due to {explain}."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"
               
        else:
            for category, df_category in filtered_dfs.items():
                sample_size = len(df_category[col])
                missing_perc = (df_category[col].isnull().sum() / sample_size) * 100
                if df_category[col].notnull().any():
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0
                    }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size: {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"

    elif mode == "tumor_filt_num":
        if exclude:
            for category, df_category in filtered_dfs.items():
                if df_category[col].notnull().any():
                    excluded_ehr = None
                    if isinstance(exclude, str):
                        excluded = df_category[col].str.contains(exclude, na=False, regex=True).sum()
                        mask = df_category[col].str.contains(exclude, na=False, regex=True)
                        df_category.loc[mask, col] = np.nan
                    else:
                        col_data = df_category[col].copy()
                        excluded_rows = col_data > exclude
                        excluded_ehr = df_category.loc[excluded_rows, 'ehr'].tolist()
                        excluded = col_data[col_data > exclude].count()
                        df_category.loc[df_category[col] > exclude, col] = np.nan

                    tumor_filter = df_category['num_tumors'] >= tumor_filt_num
                    df_category = df_category[tumor_filter]

                    sample_size = len(tumor_filter)
                    missing_perc = (df_category[col].isnull().sum() / sample_size) * 100
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc,
                        "excluded": excluded,
                        "excluded_ehr": excluded_ehr,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0,
                        "excluded": 0,
                        "excluded_ehr": 0,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }

            text_output = ''
            for category, data in sample_dict.items():
                text_output += f'{category}: Sample size: {data["sample_size"]:.2f}. Missing values: {data["missing_perc"]:.2f}%. Data excluded: {data["excluded"]:.2f} due to {explain}.'
                if category != list(sample_dict.keys())[-1]:  
                    text_output += "\n"

             
        else:
            for category, df_category in filtered_dfs.items():
                tumor_filter = df_category['num_tumors'] >= tumor_filt_num
                df_category = df_category[tumor_filter]

                sample_size = len(tumor_filter)
                missing_perc = (df_category[col].isnull().sum() / sample_size) * 100
                if df_category[col].notnull().any():
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size: {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"

    elif mode == "tumor_filt_num_no_neoadjuvant":
        if exclude:
            for category, df_category in filtered_dfs.items():
                if df_category[col].notnull().any():
                    excluded_ehr = None
                    if isinstance(exclude, str):
                        excluded = df_category[col].str.contains(exclude, na=False, regex=True).sum()
                        mask = df_category[col].str.contains(exclude, na=False, regex=True)
                        df_category.loc[mask, col] = np.nan
                    else:
                        col_data = df_category[col].copy()
                        excluded_rows = col_data > exclude
                        excluded_ehr = df_category.loc[excluded_rows, 'ehr'].tolist()
                        excluded = col_data[col_data > exclude].count()
                        df_category.loc[df_category[col] > exclude, col] = np.nan

                    tumor_filter = (df_category['num_tumors'] == tumor_filt_num) & (df_category['neoadjuvant_'+str(tumor_filt_num)] == 'no')
                    df_category_filtered = df_category[tumor_filter]

                    sample_size = len(tumor_filter)
                    missing_perc = (df_category_filtered[col].isnull().sum() / sample_size) * 100
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc,
                        "excluded": excluded,
                        "excluded_ehr": excluded_ehr,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0,
                        "excluded": 0,
                        "excluded_ehr": 0,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size (patients with {data['ordinal']} tumor no neoadjuvant): {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%. Data excluded: {data['excluded']:.2f} due to {explain}."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"

        else:
            
            for category, df_category in filtered_dfs.items():
                tumor_filter = (df_category['num_tumors'] == tumor_filt_num) & (df_category['neoadjuvant_'+str(tumor_filt_num)] == 'no')
                df_category_filtered = df_category[tumor_filter]

                sample_size = len(tumor_filter)
                missing_perc = (df_category_filtered[col].isnull().sum() / sample_size) * 100
                if df_category_filtered[col].notnull().any():
                    sample_dict[category] = {
                        "sample_size": sample_size,
                        "missing_perc": missing_perc,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }
                else:
                    sample_dict[category] = {
                        "sample_size": 0,
                        "missing_perc": 0,
                        "ordinal":get_ordinal(tumor_filt_num)
                    }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size (patients with {data['ordinal']} tumor no neoadjuvant): {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"
            
    elif mode == "tumor_filt_num_neoadjuvant":
        if exclude: 
             
            for category, df_category in filtered_dfs.items():
                if df_category[col].notnull().any():
                    excluded_ehr = None
                    if isinstance(exclude, str):
                        excluded = df_category[col].str.contains(exclude, na=False, regex=True).sum()
                        mask = df_category[col].str.contains(exclude, na=False, regex=True)
                        df_category.loc[mask, col] = np.nan
                    else:
                        col_data = df_category[col].copy()
                        excluded_rows = col_data > exclude
                        excluded_ehr = df_category.loc[excluded_rows, 'ehr'].tolist()
                        excluded = col_data[col_data > exclude].count()
                        df_category.loc[df_category[col] > exclude, col] = np.nan

                    tumor_filter = (df_category['num_tumors'] == tumor_filt_num) & (df_category['neoadjuvant_'+str(tumor_filt_num)] == 'yes')
                    df_category_filtered = df_category[tumor_filter]

                    sample_size = len(tumor_filter)
                    missing_perc = (df_category_filtered[col].isnull().sum() / sample_size) * 100
                    sample_dict[category] = {
                            "sample_size": sample_size,
                            "missing_perc": missing_perc,
                            "excluded": excluded,
                            "excluded_ehr": excluded_ehr,
                            "ordinal":get_ordinal(tumor_filt_num)
                        }
                else:
                    sample_dict[category] = {
                            "sample_size": 0,
                            "missing_perc": 0,
                            "excluded": 0,
                            "excluded_ehr": 0,
                            "ordinal":get_ordinal(tumor_filt_num)
                        }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size (patients {data['ordinal']} tumor neoadjuvant): {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%. Data excluded: {data['excluded']:.2f} due to {explain}."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"

        else: 
                
            for category, df_category in filtered_dfs.items():
                if df_category[col].notnull().any():
                        
                    tumor_filter = (df_category['num_tumors'] == tumor_filt_num) & (df_category["neoadjuvant_"+str(tumor_filt_num)] == 'yes')
                    df_category_filtered = df_category[tumor_filter]

                    sample_size = len(tumor_filter)
                    missing_perc = (df_category_filtered[col].isnull().sum() / sample_size) * 100
                    sample_dict[category] = {
                            "sample_size": sample_size,
                            "missing_perc": missing_perc,
                            "excluded": excluded,
                            "excluded_ehr": excluded_ehr,
                            "ordinal":get_ordinal(tumor_filt_num)
                        }
                else:
                        sample_dict[category] = {
                            "sample_size": 0,
                            "missing_perc": 0,
                            "excluded": 0,
                            "excluded_ehr": 0,
                            "ordinal":get_ordinal(tumor_filt_num)
                        }

            text_output = ""
            for category, data in sample_dict.items():
                text_output += f"{category}: Sample size (patients with {data['ordinal']} tumor neoadjuvant): {data['sample_size']:.2f}. Missing values: {data['missing_perc']:.2f}%."
                if category != list(sample_dict.keys())[-1]:  # Check if it is not the last line
                    text_output += "\n"


    return text_output
def numerical_variable_slide(df, prs, title_slide, title_slide_separate,  
                              divide_col, sample_string, 
                              col, exclude=None, xlim=None, ylim=None, 
                           ):
    
    df_copy = df.copy()

    if exclude: 
        excluded_mask = df_copy[col].str.contains(str(exclude), na=False, case=False)
        df_copy.loc[excluded_mask, col] = np.nan
    
    if xlim: 
        xlim_mask = df_copy[df_copy[col] > xlim]
        df_copy.loc[xlim_mask.index, col]= np.nan
    

    unique_categories = df_copy[divide_col].dropna().unique()

    filtered_dfs = {}

    for category in unique_categories:
        filtered_dfs[category] = df_copy[df_copy[divide_col] == category].copy()

    # Create figure
    fig = plt.figure(figsize=(12, 10))

   
    ax1 = fig.add_subplot(2, 2, 1)
    colors = sns.color_palette('Set2')
    for i, (group, filtered_df) in enumerate(filtered_dfs.items()):
        sns.histplot(data=filtered_df, x=col, color=colors[i], alpha=0.5, label=group, ax=ax1)
        
    ax1.legend()

    ax2 = fig.add_subplot(2, 2, 2)
    box_data = [filtered_df[col].values for filtered_df in filtered_dfs.values()]
    
    box_labels = list(filtered_dfs.keys())
    sns.boxplot(data=box_data, palette=colors, ax=ax2, orient="h")
    ax2.set_yticklabels(box_labels)

    ax3 = fig.add_subplot(2, 2, 3)
    for i, (group, filtered_df) in enumerate(filtered_dfs.items()):
        sns.histplot(data=filtered_df, x=col, kde=True, element="step", color=colors[i],  alpha=0.01, ax=ax3, edgecolor="None")
       

    if xlim:
        ax1.set_xlim([0,xlim])
        ax2.set_xlim([0, xlim])
        ax3.set_xlim([0,xlim])
    if ylim:
        ax1.set_ylim([0,ylim])
        ax3.set_ylim([0,ylim]) 

    ax1.set_xlabel(col)
    ax2.set_xlabel(col)
    ax3.set_xlabel(col)

    ax1.set_ylabel('Count')
    
    
    ax3.set_ylabel('Density')


    
    stats_dict = {}
    text_output = ""

    if len(df_copy[col].dropna()) > 3: 
        stat, p_value = shapiro(df_copy[col].dropna())
        stat, p_value = shapiro(df_copy[col].dropna())
                # Determine whether to display mean and standard deviation or median and IQR
        if p_value >= 0.05:
            normal="yes"
            f_statistic, p_value = f_oneway(*[df_copy[col][df_copy[divide_col] == cat] for cat in unique_categories])
            p = p_value if p_value >= 0.001 else "<0.001*"
        else:
            normal="no"
            p = kruskal(df_copy, dv=col, between=divide_col)["p-unc"].values[0]
                   
            if p < 0.001:
                        # Format the p-value in bold
                p = " < 0.001 ** "
            elif p < 0.05:
                p = " < 0.05 * "
                    
            else: 
                p = round(p, 4)
                


    for group_name, group_df in filtered_dfs.items():
        if group_df[col].notnull().any():
            
    # Perform Shapiro-Wilk test for normality
            if normal == "yes":
                measure_center = "Mean"
                measure_spread = "SD"
                center = statistics.mean(group_df[col].dropna())
                spread = statistics.stdev(group_df[col].dropna())
            else:
                measure_center = "Median"
                measure_spread = "IQR"
                center = statistics.median(group_df[col].dropna())
                spread = np.nanpercentile(group_df[col].dropna(), 75) - np.nanpercentile(group_df[col].dropna(), 25)
               


            if xlim:
                out_of_range_patients = group_df[(group_df[col] < 0) | (group_df[col] > xlim)]
            
            IQR=np.nanpercentile(group_df[col].dropna(), 75) - np.nanpercentile(group_df[col].dropna(), 25)
            Q1 = group_df[col].quantile(0.25)
            Q3 = group_df[col].quantile(0.75)


            stats_dict[group_name] = {
                "center": center,
                'spread': spread,
                "measure_center":measure_center, 
                "measure_spread": measure_spread, 
                'measure': spread,
                'min': group_df[col].min(),
                'max': group_df[col].max(),
                'range': group_df[col].max() - group_df[col].min(),
                'outliers': group_df[col].dropna()[~group_df[col].between(Q1 - 1.5 * IQR, Q3 + 1.5 * IQR)],
                "significance": p
            }

            if xlim:
                stats_dict[group_name]['out_of_range_ehrs'] = out_of_range_patients['ehr'].values.tolist()
                stats_dict[group_name]['out_of_range_num'] = len(out_of_range_patients)

        else:
            stats_dict[group_name] = {
                "measure":0, 
                'center': 0,
                'spread': 0,
                 "measure_center":measure_center, 
                "measure_spread": measure_spread, 
                'min': 0,
                'max': 0,
                'range': 0,
                'iqr': 0,
                'outliers': "",
                "significance": ""
            }

            if xlim:
                stats_dict[group_name]['out_of_range_ehrs'] = 0
                stats_dict[group_name]['out_of_range_num'] = 0

            text_output += f"All data null for group {group_name} for variable {col}\n"

    # Create text output
    text_output += "\nDistribution:\n"
    for group_name, group_stats in stats_dict.items():
        text_output += f"{group_name}: {group_stats['measure_center']}: {group_stats['center']:.2f}, Spread ( {group_stats['measure_spread']}): {group_stats['spread']:.2f}\n"


    text_output += f"\nData ranges:\n"
    for group_name, group_stats in stats_dict.items():
        text_output += f"{group_name}: {group_stats['min']:.2f} to {group_stats['max']:.2f}\n"

    text_output += f"\nBoxplot:\n"
    for group_name, group_stats in stats_dict.items():
        text_output += f"{group_name}: Outliers: {len(group_stats['outliers'])}\n"
    
    text_output += f"\nSignificance: {p} "
    
    
    

    if xlim:
        if not out_of_range_patients.empty:
            text_output += f"\nData out of range (0, {xlim}):\n"
            for group_name, group_stats in stats_dict.items():
                text_output += f"{group_name}: {group_stats['out_of_range_num']}\n"

      
    text_box = fig.add_subplot(2, 2, 4)
    text_box.text(0.05, 0.5, text_output, fontsize=10, verticalalignment='center', 
                color="black")

    text_box.axis('off')

    fig.savefig("plot_numerical_"+str(col)+".png")

    num_plots = len(unique_categories)
    num_cols = 2  # Number of columns in the subplot grid
    num_rows = (num_plots + num_cols - 1) // num_cols  # Calculate the number of rows needed

    colors = sns.color_palette('Set2')
    figures, axes = plt.subplots(num_rows, num_cols, figsize=(12, 8))
    figures.suptitle(f'{col} Histograms')

    axes = axes.flatten()  # Flatten the axes array for easier indexing

    for i, category in enumerate(unique_categories):
        ax = axes[i]  # Get the current subplot axis

        filtered_df = df_copy[df_copy[divide_col] == category]
        sns.histplot(data=filtered_df, x=col, color=colors[i % len(colors)], alpha=0.5, ax=ax)
        ax.set_title(f'{category}: {col} histogram')
        ax.set_xlabel(col)
        ax.set_ylabel('Count')

        if xlim:
            ax.set_xlim([0, xlim])
        if ylim:
            ax.set_ylim([0, ylim])

    # Remove any extra empty subplots
    if i < num_rows * num_cols - 1:
        for j in range(i + 1, num_rows * num_cols):
            figures.delaxes(axes[j])

    # Adjust spacing between subplots
    figures.tight_layout()
    figures.savefig("subplots_"+str(col)+".png")

    img_path = "plot_numerical_"+str(col)+".png"
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    title = slide.shapes.title
    title.left = Cm(1.27)
    title.top= Cm(0.51)
    title.height=Cm(0.56)
    title.width=Cm(22.86)
    title.text = title_slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True

    left=Cm(3.89)
    top=Cm(2.59)
    pic=slide.shapes.add_picture(img_path, left, top, height=Cm(14.67), width=Cm(17.92))

    left= Cm(0.69)
    top= Cm(1.32)
    width = Cm(24.02)
    height = Cm(2.36)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True # not necessary for newly-created shape
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    #text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = sample_string
    font = run.font
    font.size = Pt(13)


    img_path = "subplots_"+str(col)+".png"
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    title = slide.shapes.title
    title.left = Cm(1.27)
    title.top= Cm(0.51)
    title.height=Cm(0.56)
    title.width=Cm(22.86)
    title.text = title_slide_separate
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True

    left=Cm(3.53)
    top=Cm(3.87)
    pic=slide.shapes.add_picture(img_path, left, top)

from matplotlib.cm import tab10


def barchart_variable_slide(df, prs, title_slide, divide_col, sample_string, col=None, exclude=None):
    df_copy = df.copy()

    if exclude:
        excluded = len(df_copy[df_copy[col].astype(str).str.contains(str(exclude), na=False, case=False)])
        df_copy = df_copy[~df_copy[col].astype(str).str.contains(str(exclude), na=False, case=False)]

    unique_categories = df_copy[divide_col].dropna().unique()
    filtered_dfs = {}

    for category in unique_categories:
        filtered_dfs[category] = df_copy[df_copy[divide_col] == category].copy()

    category_counts = {}
    for category in unique_categories:
        category_counts[category] = filtered_dfs[category][col].astype(str).value_counts()

    # Calculate proportions for each category in each group
    category_total_counts = {}
    for category in df_copy[col].dropna().astype(str).unique():
        category_total_counts[category] = sum([counts.get(category, 0) for counts in category_counts.values()])

    proportions = {}
    for group in filtered_dfs.keys():
        proportions[group] = {}
        for category in df_copy[col].dropna().astype(str).unique():
            count = category_counts[group].get(category, 0)
            proportions[group][category] = count / category_total_counts[category] if category_total_counts[category] > 0 else 0

    categories = df_copy[col].dropna().unique()

    if np.issubdtype(df_copy[col].dtype, np.number):
        numeric_categories = pd.to_numeric(categories, errors='coerce')
        numeric_categories_sorted = np.sort(numeric_categories[~np.isnan(numeric_categories)])
        string_categories = [str(cat) for cat in categories if not np.isin(cat, numeric_categories_sorted)]
        categories_order = list(numeric_categories_sorted) + string_categories
        categories_order = list(map(str, categories_order))
    else:
        categories_order = list(map(str, categories))

    colors = sns.color_palette('Set2', n_colors=len(df[col].dropna().unique()))

    # Create figure and axes
    fig, ax = plt.subplots(figsize=(8, 6))
        # Generate a list of colors using the tab10 color palette
    num_groups = len(filtered_dfs)
    colors = tab10.colors[:num_groups]

    total_heights = np.zeros(len(categories_order))
    for i, (group, filtered_df) in enumerate(filtered_dfs.items()):
        heights = np.array([proportions[group].get(cat, 0) for cat in categories_order])
        if len(heights) == len(total_heights):
        
            ax.bar(categories_order, heights, bottom=total_heights, label=group, color=colors[i % num_groups], width=0.6)
            total_heights += heights

    # Set y-axis limit
    ax.set_ylim([0, 1])

    # Add labels and legend
    if col:
        ax.set_xlabel(col)

    ax.set_ylabel('Proportion per group')
    ax.set_title("Barchart for variable " + str(col))
    ax.legend()



    # Generate comments for each group
    comments = []
    for group, filtered_df in filtered_dfs.items():
        categories_count = len(filtered_df[col].value_counts())
        mode = filtered_df[col].dropna().mode().iloc[0] if filtered_df[col].notnull().any() else None

        comment = f"{group}: Categories: {categories_count}. Mode: {mode}."
        comments.append(comment)
    

    contingency_table = pd.crosstab(df_copy[col], df_copy[divide_col], dropna=False)
    contingency_table = contingency_table.loc[contingency_table.sum(axis=1) > 0]  # Drop categories with empty cells


    chi2, p, _, _ = chi2_contingency(contingency_table)

    if p < 0.001: 
        p = "< 0.001**"
    elif p < 0.05: 
        p = "< 0.05*"
    else: 
        p = round(p, 4)

    comments.append(f"Significance: {p}")


    # Rotate x-axis labels
    plt.xticks(rotation=45)

    # Add comments to the graph
    comment_text = "\n".join(comments)
    ax.text(0.95, 0.05, comment_text, fontsize=10, transform=ax.transAxes,
            verticalalignment='bottom', horizontalalignment='right',
            bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    
    
    # Adjust margins
    plt.subplots_adjust(bottom=0.3)

    fig.savefig("barchart_" + str(col) + ".png")

    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.left = Cm(1.48)
    title.top = Cm(1.92)
    title.height = Cm(1.12)
    title.width = Cm(22.86)
    title.text = title_slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True

    left = Cm(3.72)
    top = Cm(6.72)
    width = Cm(16.94)
    height = Cm(11.28)
    pic = slide.shapes.add_picture("barchart_" + str(col) + ".png",
                                   left, top)
                                   #, width, height)

    left = Cm(0.98)
    top = Cm(3.94)
    width = Cm(23.87)
    height = Cm(2.25)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True  # not necessary for newly-created shape
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
    p = text_frame.paragraphs[0]

    run = p.add_run()

    run.text = sample_string
    font = run.font
    font.size = Pt(14)



def intro_slide(prs, n_groups, divide_var, explain_text):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Bivariate descriptive analysis of breast cancer patients data'
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title_shape.text_frame.paragraphs[0].font.bold = True
    # add bullet point with analysis performed sentence
    analysis_bullet = slide.shapes.placeholders[1]
    analysis_bullet.text = "Data for all variables divided into "+str(n_groups)+" groups according to "+str(divide_var)

    # add text frame with categories and bullets
    category_text_frame = slide.shapes.add_textbox(left=Cm(2.22), top=Cm(8.19), 
                                                width=Cm(18.71), height=Cm(4.53))
    category_text_frame.text = explain_text


def grouper_distribution_slide(df, prs, col, colors, title_slide, exclude=None): 
    
    if exclude: 
        excluded = len(df[df[col].str.contains(str(exclude), na=False, case=False)])
        df = df[~df[col].str.contains(str(exclude), na=False, case=False)] 
      
    if colors is None:
        colors = {}

    categories = df[col].dropna().unique()
    labels = []
    counts = []

    for category in categories:
  
        if category in colors:
     
            color = colors[category]
        else:
            color = plt.cm.Set3(len(colors))
            colors[category] = color

        count = len(df[df[col] == category])

        # Add the category label and count to their respective lists
        labels.append(category)
        counts.append(count)

    # Plot the pie chart with the assigned colors and labels
    fig, ax = plt.subplots()
    wedges, _, _ = ax.pie(counts, colors=[colors[category] for category in categories], autopct='%1.1f%%', textprops={'fontsize': 10})
    ax.legend(wedges, labels, title=col, loc='center left', bbox_to_anchor=(1, 0, 0.5, 1))
    
    plt.xticks(rotation=45)
    plt.title("Piechart for grouper variable "+str(col))
    fig.savefig("piechart_grouper.png")

    fig = plt.figure()
    df.groupby(col).size().plot(kind='bar', rot=0, color="lightskyblue")
    plt.ylabel("")
    plt.title("Barchart for grouper variable "+str())
    categories = len(df[col].value_counts())
    fig.savefig("barchart_grouper.png")
   
       
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    # Add title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title_slide
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
    title_placeholder.text_frame.paragraphs[0].font.bold = True

    # Add pie chart image
    pie_chart_left = Cm(12.7)
    pie_chart_top = Cm(4.99)
    pie_chart_width = Cm(10.41)
    pie_chart_height = Cm(7.34)
    slide.shapes.add_picture("piechart_grouper.png", pie_chart_left, pie_chart_top,
                             width=pie_chart_width, height=pie_chart_height)

    # Add bar chart image
    bar_chart_left = Cm(0.86)
    bar_chart_top = Cm(4.85)
    bar_chart_width = Cm(11.26)
    bar_chart_height = Cm(7.73)
    slide.shapes.add_picture("barchart_grouper.png", bar_chart_left, bar_chart_top,
                             width=bar_chart_width, height=bar_chart_height)

    # Add textbox with explanation
    textbox_left = Cm(3.71)
    textbox_top = Cm(13.44)
    textbox_width = Cm(14.76)
    textbox_height = Cm(4.02)
    shape = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = shape.text_frame
    p = text_frame.add_paragraph()
    p.text = "Grouper variable: " + str(col) + \
         "\ncategories: " + str(df[col].dropna().unique()) + \
         "\nMode is " + str(stats.mode(df[col].dropna())[0][0]) + \
         " with " + str(stats.mode(df[col].dropna())[1][0]) + " counts." + \
         "\nMissing data: " + str(len(df[df[col].isnull()]))
    
    font = p.font
    font.name = 'Calibri'
    font.size = Pt(18)


def comorbidity_plot(df, prs, divide_col, divide_name, comorbidities): 

    value_col = divide_col
    # Define values for the value column to create subplots
    values = df[value_col].dropna().unique()
    colors = sns.color_palette('Set2', 2)
    # Calculate number of rows and columns needed for subplots
    num_plots = len(values)
    num_rows = int(num_plots / 2) + num_plots % 2
    num_cols = 2
    fig, axs = plt.subplots(num_rows, num_cols, figsize=(16, 5*num_rows))
    axs = axs.flatten()

    # Loop over each subplot and plot the stacked bar chart for each comorbidity
    for i, value in enumerate(values):
        df_sub = df[df[value_col] == value]
        si_prop = []
        no_prop = []
        for comorbidity in comorbidities:
            # Calculate proportions of si and no for each comorbidity
            si = df[comorbidity].value_counts().get('si', 0)
            no = df[comorbidity].value_counts().get('no', 0)
            total = len(df[comorbidity].dropna())
            si_prop.append(si/total * 100)
            no_prop.append(no/total * 100)

    # Create stacked bar chart

        axs[i].bar(comorbidities, si_prop, color=colors[0], label='si')
        #ax.bar(comorbidities, no_prop, bottom=si_prop, color=colors[1], label='no')
        axs[i].bar(comorbidities, no_prop, bottom=si_prop, color=colors[1], label='no')
        # Set y-axis limits to [0, 100]
        axs[i].set_ylim([0, 100])


        axs[i].set_title(f'value={value}')
        axs[i].set_ylabel('Count')
        axs[i].set_xticks(range(len(comorbidities)))
        axs[i].set_xticklabels(comorbidities, rotation=45, ha='right')
        axs[i].legend()
    # axs[i].set_ylim([0,1])
    plt.suptitle("Comorbidity count per ")
    plt.tight_layout()

    plt.savefig("comorbidities_bivariate_"+str(divide_col)+".png", bbox_inches="tight")

    used_data_comment = "Used data: "
    most_common_comment = "Most common comorbidity: "
    missing_values_comment = "Total missing values for all comorbidities: "
    excluded_values_comment= "Number of excluded rows (yes, no in same cell for any of comorbidities): "
    excluded_values = (df[comorbidities] == "no,si") | (df[comorbidities] == "si,no")
    excluded_values_count = excluded_values.any(axis=1).sum()

    for tumor_type in values:
        tumor_type_df = df[df[divide_col] == tumor_type]
        excluded_values = (tumor_type_df[comorbidities] == "si,no") | (tumor_type_df[comorbidities] == "no,si")
        # Get subset of dataframe for tumor type
        
        sample_size = len(tumor_type_df)-excluded_values_count
        used_data_comment += f"{tumor_type}: {sample_size}, "
        # Get most common comorbidity
        most_common = tumor_type_df[comorbidities].apply(pd.Series.value_counts).T.sort_values(by=['si'], ascending=False).index[0]
        most_common_comment += f"{tumor_type}: {most_common}, "
        
        # Get number of missing values
        missing_values = tumor_type_df[comorbidities].isnull().sum()
        missing_values_comment += f"{tumor_type}: {missing_values.sum()}, "
        
        # Get number of excluded values
        
        excluded_values_count = excluded_values.any(axis=1).sum()
        excluded_values_comment += f"{tumor_type}: {excluded_values_count},"

            
        # Join comments into a single string
        comments = f"{used_data_comment}\n{most_common_comment}\n{missing_values_comment}\n{excluded_values_comment}" 

    img_path = "comorbidities_bivariate_"+str(divide_col)+".png"
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)

    title = slide.shapes.title
    title.left=Cm(6.95)
    title.top=Cm(1.21)
    title.height=Cm(1.31)
    title.width=Cm(13.69)
    title.text = "Comorbidity analysis"

    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(30)


    left=Cm(2.79)
    top=Cm(2.54)
    pic=slide.shapes.add_picture(img_path, left, top, width=Cm(20.74), height=Cm(12.91))


    left= Cm(0.18)
    top= Cm(15.45)
    width = Cm(24.9)
    height = Cm(3.3)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = comments
        #text 
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(12)	


def piechart_variable_slide(df, prs, sample_string, title_slide, divide_col, colors=None, col=None, exclude=None):
    df_copy = df.copy()
    
    if exclude:
        excluded_mask = df_copy[col].str.contains(str(exclude), na=False, case=False)
        df_copy.loc[excluded_mask, col] = np.nan
    
    unique_categories = df_copy[divide_col].dropna().unique()
    filtered_dfs = {}
    
    for category in unique_categories:
        filtered_dfs[category] = df_copy[df_copy[divide_col] == category].copy()
    
    num_groups = len(filtered_dfs)
    rows = math.ceil(num_groups / 2)
    
    fig, axs = plt.subplots(rows, 2, figsize=(12, 6 * rows), squeeze=False)

    fig.suptitle("Piechart for variable " + str(col), fontsize=16)
    
    if num_groups == 1:
        
        axs = np.array([[axs[0, 0]]])

    
    for i, (group, filtered_df) in enumerate(filtered_dfs.items()):
        row_index = i // 2
        col_index = i % 2
        
        categories = filtered_df[col].dropna().unique()
        
        if len(categories) > 0:
            counts = filtered_df[col].value_counts().values
            pie = filtered_df[col].value_counts(sort=False).plot(
                kind='pie',
                textprops={'fontsize': 12},
                colors=[colors[key] for key in categories],
                autopct=lambda pct: "{:.1f}%\n({:.0f})".format(pct, pct/100.*sum(counts)),
                pctdistance=0.5,
                ax=axs[row_index, col_index] if num_groups > 1 else axs[0, 0],
            )
            axs[row_index, col_index].set_title(group)
            
            # Apply shadow effect
            for idx, wedge in enumerate(pie.patches):
                wedge.set_edgecolor('white')
                wedge.set_alpha(0.8)
        
        else:
            axs[row_index, col_index].axis('off')
    
    for ax_row in axs:
        for ax in ax_row:
            ax.set(ylabel=None)
    
    categories = [len(filtered_df[col].value_counts()) for filtered_df in filtered_dfs.values()]
    
    modes_list = []
    for filtered_df in filtered_dfs.values():
        if filtered_df[col].notnull().any():
            modes_list.append(filtered_df[col].dropna().mode().iloc[0])
        else:
            modes_list.append(None)
    
    group_names = [f"{group}: {len(filtered_df)}" for group, filtered_df in filtered_dfs.items()]
    
    comment = f"Categories: {', '.join(map(str, df[col].dropna().unique()))}\n"
    
    for group, category_count, group_mode in zip(group_names, categories, modes_list):
        comment += f"{group}: {category_count} categories, Mode: {group_mode}\n"
    
    contingency_table = pd.crosstab(df_copy[col], df_copy[divide_col], dropna=False)
    contingency_table = contingency_table.loc[contingency_table.sum(axis=1) > 0]  # Drop categories with empty cells
    chi2, p, _, _ = chi2_contingency(contingency_table)

    if p < 0.001: 
        p = "< 0.001**"
    elif p < 0.05: 
        p = "< 0.05*"
    else: 
        p = round(p, 4)

    comment += f"Significance: {p}"
    
    # Add comment below the title
    fig.text(0.5, 0.1, comment, fontsize=12, ha='center', va='top')
 
    
    fig.savefig("piechart_"+str(col)+".png", bbox_inches="tight")
    graph_slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.left = Cm(1.95)
    title.top= Cm(0.04)
    title.height=Cm(1.12)
    title.width=Cm(22.86)
    title.text = title_slide
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(22)
    title.text_frame.paragraphs[0].font.bold = True
    
    left= Cm(5.31)
    top= Cm(3.62)
    width = Cm(14.79)
    height = Cm(15.15)
    img_path ="piechart_"+str(col)+".png"
    pic=slide.shapes.add_picture(img_path, left, top)
                                 #width, height)


    left= Cm(1.27)
    top= Cm(1.16)
    width = Cm(22.86)
    height = Cm(2.4)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)

    text_frame.word_wrap = True
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    
    run = p.add_run()

    run.text = sample_string
    font = run.font
    font.size = Pt(13)

    return fig


def divide_patients_per_n_tumors(df):
    df_filtered_1 = df[df["num_tumors"] == 1]
    df_filtered_2 = df[df['num_tumors'] > 1]
    date_columns = df_filtered_2.filter(regex='tumor_date_').columns.tolist()
    df_filtered_2 = df_filtered_2[df_filtered_2[date_columns].nunique(axis=1) == 1]
    df_filtered_3 = df[df['num_tumors'] > 1]
    date_columns = df_filtered_3.filter(regex='tumor_date_').columns.tolist()
    df_filtered_3 = df_filtered_3[df_filtered_3[date_columns].nunique(axis=1) > 1]
        
    df_combined = pd.concat([df_filtered_1, df_filtered_3])
    
    return df_combined, df_filtered_2



def plot_surgeries(df, prs,  divide_col): 

    divide_col_categories= df[divide_col].dropna().unique()
    surgery_cats = ["lymphadenectomy", "partial mastectomy", "sentinel lymph node biopsy", "mastectomy"]

    # Create a dictionary to store the surgery counts for each divide_col category
    counts = {}
    filtered_dfs = {}
    for category in divide_col_categories:
        counts[category] = {}
        filtered_df = df[df[divide_col] == category]
        filtered_dfs[category] = filtered_df
        surgery_counts = filtered_df[surgery_cats].sum().to_dict()
        for surgery_type, count in surgery_counts.items():
            counts[category][surgery_type] = count

    proportions = {}
    for surgery_type in surgery_cats:
        proportions[surgery_type] = {}
        total_count = sum(counts[category][surgery_type] for category in divide_col_categories)
        for category in divide_col_categories:
            proportions[surgery_type][category] = counts[category][surgery_type] / total_count

    categories = list(proportions.keys())
    data = [[proportions[c][g] for c in categories] for g in divide_col_categories]

    # Define colors based on the number of categories
    my_palette=sns.color_palette('Set2', 8)

    # Create a stacked horizontal bar chart
    fig, ax = plt.subplots()
    for i in range(len(divide_col_categories)):
        ax.barh(categories, data[i], left=np.sum(data[:i], axis=0), color=my_palette[i], label=divide_col_categories[i])

    # Split the y-tick labels into two lines if there is a space between words
    new_labels = []
    for label in categories:
        if ' ' in label:
            idx = label.find(' ')
            new_label = label[:idx] + '\n' + label[idx+1:]
            new_labels.append(new_label)
        else:
            new_labels.append(label)

    ax.set_yticks(np.arange(len(new_labels)))
    ax.set_yticklabels(new_labels)

    # Add axis labels and legend
    ax.set_ylabel('Types of surgeries')
    ax.set_xlim([0, 1])
    ax.set_xlabel('Proportions')
    ax.legend()

    # Generate comments for each group
    comments = []
    for group, filtered_df in filtered_dfs.items():
        surgery_columns = filtered_df[surgery_cats]
        mode_column = surgery_columns.sum().idxmax() if not surgery_columns.empty else None
        comment = f"{group}: Mode: {mode_column}."
        comments.append(comment)

    # Rotate x-axis labels
    plt.xticks(rotation=45)

    # Add comments to the graph
    comment_text = "\n".join(comments)
    ax.text(0.35, 0.05, comment_text, fontsize=10, transform=ax.transAxes,
            verticalalignment='bottom', horizontalalignment='right',
            bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))

    # Adjust margins
    plt.subplots_adjust(bottom=0.3)

    # Show the chart
    plt.savefig("surgery_perc_divide_col.png")
    sample_sizes = {}
    for group in filtered_dfs.keys():
        group_df = filtered_dfs[group]
        sample_sizes[group] = len(group_df)

    #sample_string = "ample size (patients treated by surgery): " + str(sum(sample_sizes.values())) + ",\n Number of missing values: 0.0%\n"
    sample_string = "- Sample sizes: "
    missing_string = "- Missing values: "
    for group, sample_size in sample_sizes.items():
        sample_string += f"{group}: {sample_size}, "
        missing_string += f"{group}: 0.0%, "
    sample_string = sample_string[:-2] + "\n"
    missing_string = missing_string[:-2] + "\n"
    sample_string += missing_string

    fig.savefig("surgery_perc_divide_col.png")

    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.left = Cm(1.48)
    title.top= Cm(1.92)
    title.height=Cm(1.12)
    title.width=Cm(22.86)
    title.text = "Treatment analysis: percentage of patients treated by each surgery"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True

    left= Cm(3.72)
    top= Cm(6.72)
    width = Cm(16.94)
    height = Cm(11.28)
    pic=slide.shapes.add_picture("surgery_perc_divide_col.png", 
                                 left, top, width, height)

    
    
    left= Cm(0.98)
    top= Cm(3.94)
    width = Cm(23.87)
    height = Cm(2.25)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True # not necessary for newly-created shape
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    p = text_frame.paragraphs[0]

    run = p.add_run()

    run.text = sample_string
    font = run.font
    font.size = Pt(14)
        
   

def plot_drugs(df, prs,  divide_col): 

    divide_col_categories = df[divide_col].dropna().unique()
    drug_cats = ['everolimus', 'olaparib', 'exemestane', 'tamoxifen', 'alpelisib',
                  'fulvestrant', 'abemaciclib', 'letrozole', 'ribociclib', 'palbociclib', 'goserelin', 
                  'anastrozole', 'vinorelbine', 'megestrol acetate', 'capecitabine']

    # Create a dictionary to store the drug counts for each divide_col category
    counts = {}
    filtered_dfs = {}
    for category in divide_col_categories:
        counts[category] = {}
        filtered_df = df[df[divide_col] == category]
        filtered_dfs[category] = filtered_df
        drug_counts = filtered_df[drug_cats].sum().to_dict()
        for drug, count in drug_counts.items():
            counts[category][drug] = count

    proportions = {}
    for drug in drug_cats:
        proportions[drug] = {}
        total_count = sum(counts[category][drug] for category in divide_col_categories)
        for category in divide_col_categories:
            if total_count != 0: 
                proportions[drug][category] = counts[category][drug] / total_count  
            else: 
                proportions[drug][category] = 0 

    categories = list(proportions.keys())
    data = [[proportions[d][g] for d in drug_cats] for g in divide_col_categories]


    my_palette=sns.color_palette('Set2', 12)
    # Create a stacked horizontal bar chart
    fig, ax = plt.subplots()
    for i in range(len(divide_col_categories)):
        ax.barh(drug_cats, data[i], left=np.sum(data[:i], axis=0), color=my_palette[i], label=divide_col_categories[i])

    # Split the y-tick labels into two lines if there is a space between words
    new_labels = []
    for label in drug_cats:
        if ' ' in label:
            idx = label.find(' ')
            new_label = label[:idx] + '\n' + label[idx+1:]
            new_labels.append(new_label)
        else:
            new_labels.append(label)

    ax.set_yticks(np.arange(len(new_labels)))
    ax.set_yticklabels(new_labels)

    # Add axis labels and legend
    ax.set_ylabel('Drug Categories')
    ax.set_xlim([0, 1])
    ax.set_xlabel('Proportions')
    ax.legend()

    # Generate comments for each group
    comments = []
    for group, filtered_df in filtered_dfs.items():
        drug_columns = filtered_df[drug_cats]
        mode_column = drug_columns.sum().idxmax() if not drug_columns.empty else None
        comment = f"{group}: Mode: {mode_column}."
        comments.append(comment)

    # Rotate x-axis labels
    plt.xticks(rotation=45)
    plt.yticks(fontsize=9)
    # Add comments to the graph
    comment_text = "\n".join(comments)
    ax.text(0.35, 0.05, comment_text, fontsize=10, transform=ax.transAxes,
            verticalalignment='bottom', horizontalalignment='right',
            bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))

    # Adjust margins
    plt.subplots_adjust(bottom=0.3)


    sample_sizes = {}
    for group in filtered_dfs.keys():
        group_df = filtered_dfs[group]
        sample_sizes[group] = len(group_df)
    sample_string = "- Sample sizes: "
    missing_string = "- Missing values: "
    for group, sample_size in sample_sizes.items():
        sample_string += f"{group}: {sample_size}, "
        missing_string += f"{group}: 0.0%, "
    sample_string = sample_string[:-2] + "\n"
    missing_string = missing_string[:-2] + "\n"
    sample_string += missing_string

    # Show the chart
    plt.savefig("drug_perc_divide_col.png")
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.left = Cm(1.48)
    title.top= Cm(1.92)
    title.height=Cm(1.12)
    title.width=Cm(22.86)
    title.text = "Treatment analysis: percentage of patients treated by each Ht drug"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
    title.text_frame.paragraphs[0].font.size = Pt(25)
    title.text_frame.paragraphs[0].font.bold = True

    left= Cm(3.72)
    top= Cm(6.72)
    width = Cm(16.94)
    height = Cm(11.28)
    pic=slide.shapes.add_picture("drug_perc_divide_col.png", 
                                 left, top, width, height)



def descriptive_summary_definitiva(df, group_col):
    if group_col not in df.columns:
        raise ValueError(f"{group_col} is not a valid column in the DataFrame.")

    df_copy = df.copy()
    unique_groups = df_copy[group_col].dropna().unique()
    summary_df = pd.DataFrame(columns=['Variable', 'Category', 'N', *unique_groups, 'p-value'])

    # Iterate over each column in the dataframe
    for column in df_copy.columns:
        if column != group_col:
            data = df_copy[[group_col, column]].copy()
            data[column] = data[column].replace("Exclude", np.nan)

            # Categorical variables
            if data[column].dtype in ['object', 'category']:
                contingency_table = pd.crosstab(data[column], data[group_col], dropna=False)
                contingency_table = contingency_table.loc[contingency_table.sum(axis=1) > 0]  # Drop categories with empty cells

                if contingency_table.empty:
                    continue  # Skip this column if all categories have empty cells

                chi2, p, _, _ = chi2_contingency(contingency_table)

              

                # Add a row for the categorical variable with only variable name and total counts
                variable_row = {
                    'Variable': column,
                    'Category': '',
                    'N': f'N={len(df_copy[column])}',
                    'p-value': "",
                }
                  
                
                if p < 0.001:
                        # Format the p-value in bold
                    variable_row['p-value'] = " < 0.001 ** "
                elif p < 0.05:
                    variable_row['p-value'] = " < 0.05 * "
                    
                else: 
                    p = round(p, 4)
                    variable_row['p-value'] = p





                summary_df = pd.concat([summary_df, pd.DataFrame([variable_row])], ignore_index=True)

                for category in contingency_table.index:
                    total_count = len(df_copy[df_copy[column] == category])

                    result = {
                        'Variable': '',
                        'Category': category,
                        'N': f'N={total_count}',
                        'p-value': '',
                    }

                    for group in unique_groups:
                        count = contingency_table.loc[category, group]
                        total_count = contingency_table[group].sum()
                        percentage = (count / total_count) * 100
                        result[group] = f'{count} ({percentage:.2f}%)'

                    summary_df = pd.concat([summary_df, pd.DataFrame([result])], ignore_index=True)
            else:
                result = {
                    'Variable': column,
                    'Category': '',
                    'N': f'N={len(data)}',
                    'p-value': '',
                }
                _, p_normality = shapiro(df_copy[column].dropna())
                
                if p_normality < 0.05:
                        # Variable is not normal, perform non-parametric test (Kruskal-Wallis)
                    
                    
                    for group in unique_groups:
                        group_data = data[data[group_col] == group][column]
                        median = group_data.median()
                        q1 = group_data.quantile(0.25)
                        q3 = group_data.quantile(0.75)
                        iqr = q3 - q1
                        
                        result[group] = f'{median:.2f}±{iqr:.2f}'
                    
                    p = kruskal(df_copy, dv=column, between=group_col)["p-unc"].values[0]
                   
                    if p < 0.001:
                        # Format the p-value in bold
                        result['p-value'] = " < 0.001 ** "
                    elif p < 0.05:
                        result['p-value'] = " < 0.05 * "
                    
                    else: 
                        p = round(p, 4)
                        result['p-value'] = p


                else:

                    for group in unique_groups:
                        group_data = data[data[group_col] == group][column]
                        mean = group_data.mean()
                        sd = group_data.std()

                        result[group] = f'{mean:.2f} ± {sd:.2f}'

                    # Check for normality using Shapiro-Wilk test
                   


                        # Perform ANOVA
                    f_statistic, p = f_oneway(*[data[data[group_col] == g][column] for g in unique_groups])  # Updated indexing here
                    if p < 0.05:
                        # Format the p-value in bold
                        result['p-value'] = f'<b>{p:.4f}</b>'
                    else:
                        result['p-value'] = f'{p:.4f}'






                    if p < 0.001:
                        p = "<0.001"
                    else:
                        p = round(p, 4)
                   
                    result['p-value'] = p
                        # Variable is normal, perform ANOVA
                   

              

                summary_df = pd.concat([summary_df, pd.DataFrame([result])], ignore_index=True)

    # Round the numerical values to two decimal places
    numerical_columns = unique_groups.tolist()
    summary_df[numerical_columns] = summary_df[numerical_columns].round(3)

    # Replace NaN with empty string
    summary_df = summary_df.replace(np.nan, '', regex=True)

    return summary_df

