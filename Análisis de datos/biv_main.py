import time
start_time = time.time()


import datetime as dt
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics as stats
from pptx import Presentation
from sklearn.compose import make_column_selector as selector
from pptx.util import Inches, Pt
from pptx.util import Cm
import io
import re
from PIL import Image
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE 
import inflect
import os
import sys
import procesamiento
from procesamiento import obtain_cols
from procesamiento import to_date
from procesamiento import extract_intervals_2025
from procesamiento import any_treatments
from procesamiento import create_treatment_columns
import bivariate_tools
from bivariate_tools import intro_slide
from bivariate_tools import numerical_variable_slide
from bivariate_tools import barchart_variable_slide
from bivariate_tools import sample_string_function
from bivariate_tools import piechart_variable_slide
from bivariate_tools import divide_patients_per_n_tumors
from bivariate_tools import plot_surgeries
from bivariate_tools import plot_drugs
from bivariate_tools import grouper_distribution_slide
from bivariate_tools import comorbidity_plot
import argparse


parser = argparse.ArgumentParser(description='Este script creará un análisis univariante de sus datos en formato .csv que se salvará en un PowerPoint.')
parser.add_argument('-a', '--argument1', type=str, help='Ruta o nombre del archivo .csv para su análisis univariante.')
parser.add_argument('-b', '--argument2', type=str, help='Ruta o nombre del archivo donde se desee guardar la presentación .pptx.')
parser.add_argument('-c', '--argument3', type=str, help='Columna divisora de los datos')
args = parser.parse_args()

# Access the argument values
argument1_value = args.argument1
argument2_value = args.argument2
divide_var = args.argument3


prs=Presentation()
df = pd.read_csv(argument1_value, sep=';')

df = obtain_cols(df)
df = to_date(df)
df=extract_intervals_2025(df)
df=any_treatments(df)
df=create_treatment_columns(df, r'surgery_[0-9]+')
df=create_treatment_columns(df, r'drug_[0-9]+')


divide_name = "recurrence in cancer"
palette = sns.color_palette('Set2')
unique_categories = df[divide_var].dropna().unique()
divide_var_colors = dict(zip(unique_categories, palette))


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Bivariate descriptive analysis of breast cancer patients data'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Analysis performed for: ' + str(df.shape[0]) + " patients based on "+str(divide_name)
p = tf.add_paragraph()
p.text = 'Comorbidity analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Demography analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Ginecological analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Tumor type analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Stage analysis'
p.level = 1

p = tf.add_paragraph()
p.text = 'Time intervals analysis'
p.level = 1

intro_slide(n_groups=4, prs=prs,  divide_var="inmunohistochemical tumor type", explain_text="Based on the values of ER, PR and HER2\n\n" + \
                            "- (PR positive or ER positive) and HER2 positive: PP\n\n" + \
                            "- (PR positive or ER positive) and HER2 negative: PN\n\n" + \
                           "- (PR negative and ER negative) and HER2 positive: NP\n\n" + \
                           "- (PR negative and ER negative) and HER2 negative: NN"
  )



grouper_distribution_slide(df=df, 
                           prs=prs, 
                           col=divide_var, 
                           colors=divide_var_colors, 
                           title_slide="Distribution of the grouper variable "+str(divide_name))


comorbidity_plot(df=df, prs=prs, divide_name=divide_name, 
                 comorbidities=['autoimmune disease',
 'cardiac insufficiency',
 'diabetes',
 'dislipemia',
 'ex-smoker',
 'gastrointestinal disease',
 'hta',
 'insomnia',
 'ischemic cardiopathology',
 'liver disease',
 'lung disease',
 'musculoskeletal disease',
 'other cardiopathology',
 'psychiatric disorder',
 'renal disease',
 'smoker',
 'thyroid disease',
 'transplant'],  
                 divide_col=divide_var 
                 )

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Demography analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Age at diagnosis'
p.level = 1

numerical_variable_slide(df=df, prs=prs, title_slide="Demography analysis: age at diagnosis",
                        title_slide_separate="Separate histograms for variable age at diagnosis",  
                        divide_col = divide_var, 
                        sample_string=sample_string_function(df, mode="normal",
                                                              divide_col=divide_var , col="age_at_diagnosis"), 
                        
                          exclude=None, xlim=100, ylim=260, col="age_at_diagnosis"
                           )


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Gynecological analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Menarche age'
p.level = 1
p = tf.add_paragraph()
p.text = 'Menopause rate at diagnosis'
p.level = 1
p = tf.add_paragraph()
p.text = 'Number of births'
p.level = 1
p = tf.add_paragraph()

p.text = 'Number of abortions'
p.level = 1



numerical_variable_slide(df, prs, title_slide="Ginecological analysis: menarche age",
                        title_slide_separate="Separate histograms for variable menarche age",  
                        divide_col = divide_var, 
                        sample_string=sample_string_function(df, mode="normal",
                                                              divide_col=divide_var, 
                                                              col="menarche_age"), 
                        col="menarche_age",
                          exclude=None, xlim=20, ylim=400, 
                           )
 

piechart_variable_slide(df, prs, title_slide="Ginecological analysis: menopause at diagnosis",
                             divide_col=divide_var, colors={"yes":"lightpink", "no":"lightblue"} , 
                             sample_string=sample_string_function(df, mode="normal", 
                                                         divide_col=divide_var, col="menopause"), 
                             col="menopause", exclude=None)
    
title_slides =["Ginecological analysis: number of pregnancies in patients", 
               "Ginecological analysis: number of natural births in patients", 
               "Ginecological analysis: number of abortions in patients", 
               "Ginecological analysis: number of caesareans in patients"]
ginecological_vars =["pregnancy", "natural_birth", "abort", "caesarean"]


for gin_var, title_slide in  zip (ginecological_vars, title_slides): 
    barchart_variable_slide(df, prs, title_slide=title_slide, 
                                divide_col=divide_var,
                                sample_string=sample_string_function(df=df, mode="normal", 
                                                         divide_col=divide_var, col=gin_var), 
                                col=gin_var, exclude=None)



df_filt_1 = divide_patients_per_n_tumors(df)[0]
df_filt_2 = divide_patients_per_n_tumors(df)[1]

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True

 # assuming the body text is in the second placeholder
tf = body_shape.text_frame


tf.text = 'Variables analysed (except num. of tumors) are divided according to 3 groups of patients:'

p = tf.add_paragraph()
p.text = 'Group 1: Patients with a single tumor '+" ("+str(len(df_filt_1))+ " patients)"
p.level = 1
p.font.size = Pt(20)


p = tf.add_paragraph()
p.text = 'Group 2: Patients with more than one tumor at diagnosis '+" ("+str(len(df_filt_2))+ " patients)"
p.level = 1
p.font.size = Pt(20)


p = tf.add_paragraph()
p.text = "Only group 1 of patients where represented for this analysis (not enough data to extract significant conclusions from the rest of the patients)"
p.level = 1
p.font.size = Pt(20)



bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Nº of tumors'
p.level = 1

p = tf.add_paragraph()
p.text = 'Histological grade'
p.level = 1

p = tf.add_paragraph()

p.text = 'Histological type'
p.level = 1

p = tf.add_paragraph()
p.text = 'Behavior (in situ vs. invasive and associated in situ)'
p.level = 1


titles_slides = ["Tumor analysis: histological grade of tumor 1 (for patients with 1 tumor)",
                 "Tumor analysis: histological type of tumor 1 (for patients with 1 tumor)",
                 "Tumor analysis: behavior of tumor 1 (for patients with 1 tumor)"]


tumor_vars = ["grade_1", "histological_type_1", "behavior_1"]

for tumor_var, title_slide in zip(tumor_vars, titles_slides):
    if tumor_var in ["histological_type_1", "behavior_1"]:
        to_exclude = "Exclude"
        to_explain = " due to ambiguous value"
    else:
        to_exclude = None
        to_explain = None
    

    unique_categories = df[tumor_var].dropna().unique()
    palette = sns.color_palette('Set2')
    category_colors = dict(zip(unique_categories, palette))


    piechart_variable_slide(df_filt_1, prs, title_slide=title_slide, colors=category_colors,
                            divide_col=divide_var,
                            col=tumor_var, exclude=to_exclude, 
                            sample_string=sample_string_function(df_filt_1,
                                                                 mode="normal",
                                                                 divide_col=divide_var,
                                                                 exclude=to_exclude,
                                                                 explain=to_explain, col=tumor_var))


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Tumor analysis: staging'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Stage at diagnosis'
p.level = 1
p = tf.add_paragraph()
p.text = 'Neoadjuvance'
p.level = 1
p = tf.add_paragraph()
p.text = 'Stage after neoadjuvant (only patients with neoadjuvant)'
p.level = 1


unique_categories = df["stage_diagnosis_summ_1"].dropna().unique()
palette = sns.color_palette('Set2')
category_colors = dict(zip(unique_categories, palette))

piechart_variable_slide(df_filt_1, prs, title_slide="Staging analysis: stage of tumor 1 at diagnosis (patients with one tumor)",
                            divide_col=divide_var, exclude="x", 
                             col="stage_diagnosis_summ_1", colors=category_colors, 
                             sample_string=sample_string_function(df_filt_1, col="stage_diagnosis_summ_1", 
                                                                  mode="normal", 
                                                                   divide_col=divide_var, 
                                                                exclude="x", 
                                                                     explain=" due to unknown T, N or M"))

piechart_variable_slide(df_filt_1, prs, title_slide="Tumors treated with neoadjuvant therapy (patients with one tumor)",
                            divide_col=divide_var, colors={"yes": "lightpink", "no":"lightblue"}, 
                             col="neoadjuvant_1", 
                             sample_string=sample_string_function(df_filt_1, col="neoadjuvant_1", 
                                                                  mode="normal", 
                                                                   divide_col=divide_var 
                                                                 ))

piechart_variable_slide(df_filt_1, prs, title_slide="Staging analysis: Stage after neoadjuvant therapy for patients treated by neoadjuvant",
                            divide_col=divide_var, exclude="x",  
                             col="stage_after_neo_summ_1", colors=category_colors, 
                             sample_string=sample_string_function(df_filt_1, col="stage_after_neo_summ_1",
                                                                  mode="tumor_filt_num_neoadjuvant", 
                                                                   divide_col=divide_var, 
                                                                 exclude="x", 
                                                                     explain=" due to unknown T, N or M",
                                                                       tumor_filt_num=1))


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tumor type analysis: celular markers'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'
p = tf.add_paragraph()
p.text = 'Ki67 percentage'
p.level = 1

numerical_variable_slide(df_filt_1, prs, title_slide="Tumor marker analysis: ki67 protein level in tumors (patients with one tumor)", 
                          divide_col=divide_var, title_slide_separate="Separate histograms for variable ki67 in tumor 1", 
                            col="ki67_1", 
                            xlim=100, ylim=500, 
                            sample_string=sample_string_function(df, 
                                                                 divide_col=divide_var, 
                                                                     mode="normal", 
                                                                     col="ki67_1" ) )


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Treatment analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Patients treated by each treatment (hormonetherapy, surgery, radiotherapy)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Types of surgeries'
p.level = 1

p = tf.add_paragraph()
p.text = 'Types of drugs used in hormonetherapy'
p.level = 1

p = tf.add_paragraph()
p.text = 'Number of each type of treatment (surgery, radiotherapy, hormonetherapy)'
p.level = 1
p = tf.add_paragraph()



piechart_variable_slide(df, prs, title_slide="Treatment analysis: surgery rate in patients", colors={"yes": "lightpink", "no":"lightblue"},
                             divide_col=divide_var, 
                             sample_string=sample_string_function(df, divide_col=divide_var, col="Any_surgery", mode="normal"), 
                             col="Any_surgery", exclude=None )
piechart_variable_slide(df, prs, title_slide="Treatment analysis: Rx rate in patients",
                             divide_col=divide_var, colors={"yes": "lightpink", "no":"lightblue"},
                             sample_string=sample_string_function(df, divide_col=divide_var, col="Any_radiotherapy", mode="normal"), 
                             col="Any_radiotherapy", exclude=None )

piechart_variable_slide(df, prs, title_slide="Treatment analysis: Ht rate in patients",
                             divide_col=divide_var, colors={"yes": "lightpink", "no":"lightblue"},
                             sample_string=sample_string_function(df, divide_col=divide_var, col="Any_hormonetherapy", mode="normal"), 
                             col="Any_hormonetherapy", exclude=None )

plot_surgeries(df, prs, divide_col=divide_var)
plot_drugs(df, prs, divide_col=divide_var)



plot_surgeries(df, prs, divide_col=divide_var)
plot_drugs(df, prs, divide_col=divide_var)


barchart_variable_slide(df, prs, title_slide="Treatment analysis: number of surgeries in patients",
                             divide_col=divide_var, 
                             sample_string=sample_string_function(df, divide_col=divide_var, col="N_surgeries", mode="normal"), 
                             col="N_surgeries",  exclude=None )

barchart_variable_slide(df, prs, title_slide="Treatment analysis: number of Rx in patients",
                             divide_col=divide_var, 
                             sample_string=sample_string_function(df, divide_col=divide_var, col="n_radio", mode="normal"), 
                             col="n_radio",  exclude=None )


barchart_variable_slide(df, prs, title_slide="Treatment analysis: number of Ht in patients",
                             divide_col=divide_var, 
                             sample_string=sample_string_function(df, divide_col=divide_var, col="N_hormonetherapies", mode="normal"), 
                             col="N_hormonetherapies",  exclude=None )



bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Time interval analysis'
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
title_shape.text_frame.paragraphs[0].font.bold = True
tf = body_shape.text_frame
tf.text = 'Variables analysed'

p = tf.add_paragraph()
p.text = 'Time from dx to first treatment in all patients (Days)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Time from dx to first surgery in patients with no neoadjuvant therapy for 1st tumor (Days)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Time from dx to recurrence (Months)'
p.level = 1



numerical_variable_slide(df, prs=prs, divide_col=divide_var, title_slide="Time interval analysis: time from dx to first surgery (non neoadjuvant patients)",
                          title_slide_separate="Separate histograms for variable time to first surgery (non neoadjuvant patients)", 
                         col="Time_dx_surgery_no_neo_days_1", xlim=90, 
                         sample_string = sample_string_function(df,divide_col=divide_var, mode="tumor_filt_num",
                                                                                   col="Time_dx_surgery_no_neo_days_1",
                                                                                     exclude=90, tumor_filt_num=1))

numerical_variable_slide(df, prs=prs, divide_col=divide_var, title_slide="Time interval analysis: time from diagnosis to neoadjuvant therapy (neoadjuvant patients)",
                          title_slide_separate="Separate histograms for variable time to neoadjuvant therapy", 
                         col="Time_dx_neo_days_1", xlim=120, 
                         sample_string = sample_string_function(df,divide_col=divide_var, mode="tumor_filt_num",
                                                                                   col="Time_dx_neo_days_1",
                                                                                     exclude=120, tumor_filt_num=1))


numerical_variable_slide(df[df["recurrence"] == "yes"], prs, divide_col=divide_var, title_slide="Time interval analysis: time to recurrence in Months (patients with recurrence)",
                          title_slide_separate="Separate histograms for variable time to recurrence (Months)", 
                         col="Time_to_recurrence_months", 
                         sample_string = sample_string_function(df[df["recurrence"] == "yes"],
                                                                divide_col=divide_var, mode="normal",
                                                                                   col="Time_to_recurrence_months"
                                                                                    ))


prs.save(argument2_value+".pptx")


end_time = time.time()
execution_time = end_time - start_time

print("Script executed in {:.2f} seconds.".format(execution_time))

