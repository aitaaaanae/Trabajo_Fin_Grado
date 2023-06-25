import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from lifelines import KaplanMeierFitter
from lifelines.statistics import multivariate_logrank_test
from matplotlib.offsetbox import AnchoredText
import tempfile
import os
import getpass
import sys




def km_plot(df, prs, event_col, time_col, event_name, time_unit, group_col,  exclude=None,
            cutoff_points=None, time_cut=None):
    df_copy = df.copy()

    if exclude:
        excluded_mask = df_copy[group_col].str.contains(str(exclude), na=False, case=False)
        df_copy.loc[excluded_mask, group_col] = np.nan

    df_copy["Event"] = np.where(df_copy[event_col].str.strip().str.lower() == "yes", 1, 0)
    #censor_date = pd.to_datetime("2022-05-31")
   
    df_copy["Time"] = df_copy[time_col] //30
    df_copy = df_copy.dropna(subset=["Time"])
 
    #df_copy["Time"] = time_col
    #df_copy["Time"].fillna((pd.to_datetime(df_copy[""]) - pd.to_datetime(df_copy["diagnosis_date"])).dt.days // 30)
    #df_copy = df_copy.dropna(subset=["diagnosis_date"])

    unique_groups = df_copy[group_col].dropna().unique()
    if cutoff_points:
        cutoff_points = list(map(float, cutoff_points.strip('[]').split(',')))
        df_copy[group_col] = pd.cut(pd.to_numeric(df_copy[group_col]), cutoff_points)
        unique_groups = df_copy[group_col].dropna().unique()
        df_copy = df_copy.dropna(subset=[group_col])

    groups = [df_copy[group_col] == group for group in unique_groups]
    T = df_copy["Time"].values
    E = df_copy["Event"].values
    labels = [*unique_groups]

    num_groups = len(df_copy[group_col].unique())
    my_palette = sns.color_palette("Set2", n_colors=num_groups)

    fig, ax = plt.subplots(figsize=(10, 6))
    kmf = KaplanMeierFitter()

    for i, group in enumerate(groups):
        T_group = T[group]
        E_group = E[group]
        label = labels[i]

        kmf.fit(T_group, E_group, label=label)
        kmf.plot(ci_show=False, color=my_palette[i], ax=ax)

    ax.set_xlabel('Time (' + str(time_unit) + ') after diagnosis')
    ax.set_ylabel(str(event_name) + ' probability')
    ax.set_ylim([0, 1.1])
    ax.set_title("Kaplan Meier curves for variable "+str(group_col))
    ax.legend()


    # Perform log-rank test
    result = multivariate_logrank_test(df_copy['Time'], df_copy[group_col], df_copy['Event'])
    result.test_statistic
    result.p_value

    # Add log-rank p-value to the plot
    anchored_text = AnchoredText("log-rank p-value: {:.3}".format(result.p_value), loc='upper left', frameon=False)
    plt.gca().add_artist(anchored_text)

    if time_cut is not None and time_cut != '':
        time_cut_months = int(time_cut) * 12  # Convert years to months

        kmf.fit(T, E)
        kmf_survival_at_time_cut = kmf.predict(time_cut_months)

        ax.axvline(x=time_cut_months, color='r', linestyle='--')
        ax.text(time_cut_months + 10, 0.5, f"{event_name} at {time_cut} years", rotation=0, ha='left', va='center')
        #lt.axvline(x=time_cut_months, color='r', linestyle='--')
        #plt.text(time_cut_months + 10, 0.5, f"{event_name} at {time_cut} years", rotation=0, ha='left', va='center')
        
        for i, group in enumerate(groups):
            T_group = T[group]
            E_group = E[group]
            label = labels[i]
            
            kmf.fit(T_group, E_group, label=label)
            kmf.plot(ci_show=False, color=my_palette[i])
            ax.set_xlabel('Time (' + str(time_unit) + ') after diagnosis')
            
            kmf_survival_at_time_cut_group = kmf.predict(time_cut_months)
            
            plt.text(time_cut_months + 10, 0.5 - (i+1)*0.05, f"{label}: {kmf_survival_at_time_cut_group:.2f}", rotation=0, ha='left', va='center')

    #with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
    plt.savefig(str(group_col)+".png")
    img_path = str(group_col)+".png"

    graph_slide_layout = prs.slide_layouts[5]
   
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Kaplan Meier Curves for "+str(group_col)
    title.text_frame.paragraphs[0].font.bold = True

    slide.shapes.add_picture(img_path, left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(5))


    os.remove(img_path)
   

import streamlit as st
import pandas as pd
import os
from pptx import Presentation

def main():
    st.set_page_config(
        page_title="Survival Curve Plotter",
        page_icon=":bar_chart:",
        layout="wide",
        initial_sidebar_state="expanded",
    )

    st.markdown(
        """
        <style>
        body {
            background-color: #f2f2f2;
        }
        .stButton button {
            background-color: #FF1493;
            color: white;
        }
        .reportview-container .main .block-container {
            max-width: 800px;
        }
        </style>
       
        """,
        unsafe_allow_html=True,
    )

    st.title("Survival Curve Plotter")



    # Upload CSV file
    uploaded_file = st.file_uploader("Upload CSV file", type="csv")

    if uploaded_file is not None:
        df = pd.read_csv(uploaded_file)

        st.header("Survival Curve Parameters")

        # Create a two-column layout
        col1, col2 = st.columns(2)

        # Column 1: Event column, Time column, Event name, Group column
        with col1:
            st.markdown("## Mandatory Parameters")
            event_col = st.text_input("Event column", value="recurrence")
            time_col = st.text_input("Time column", value="Time_to_recurrence_months")
            event_name = st.text_input("Event name", value="Survival Event")
            group_col = st.text_input("Group column", value="grade_1")

        # Column 2: Time unit, Cutoff points, Time cut, Exclude
        with col2:
            st.markdown("## Optional Parameters")
            time_unit = st.text_input("Time unit", value="months")
            cutoff_points = st.text_input("Cutoff points (optional)")
            time_cut = st.text_input("Time cut (optional)")
            exclude = st.text_input("Exclude (optional)")

        # Plot and Save buttons
        st.header("Actions")
        plot_button = st.button("Plot")
        save_button = st.button("Save as PowerPoint")

        if plot_button:
            if os.path.exists("tmp.pptx"):
                prs = Presentation("tmp.pptx")
            else:
                prs = Presentation()

            km_plot(df, prs, event_col, time_col, event_name, time_unit, group_col, exclude, cutoff_points, time_cut)

            st.pyplot()
            prs.save("tmp.pptx")

        if save_button:
            ppt_file = st.text_input("Enter the file path to save the PowerPoint file:", value="plot.pptx")
            if os.path.exists(ppt_file):
                os.remove(ppt_file)
            os.rename("tmp.pptx", ppt_file)


if __name__ == "__main__":
    main()
